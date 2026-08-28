"""Strict offline contract for one candidate-bound classroom export set."""

from __future__ import annotations

from collections.abc import Mapping
from datetime import datetime
import hashlib
from html.parser import HTMLParser
import json
import os
from pathlib import Path, PurePosixPath, PureWindowsPath
import posixpath
import re
import stat
import struct
from typing import BinaryIO
from urllib.parse import quote, urlsplit
import xml.etree.ElementTree as ET
import zipfile

from pptx import Presentation
from pptx.exc import PythonPptxError
from pydantic import ValidationError

from deeptutor.teaching.contracts import ClassroomDocument, canonical_json_bytes

CLASSROOM_EXPORT_SCHEMA_VERSION = 1
CLASSROOM_EXPORT_PRODUCER = "classroom-export-probe"
CLASSROOM_EXPORT_KINDS = (
    "classroom_zip",
    "pptx",
    "offline_html",
    "mp4",
)
CLASSROOM_EXPORT_PATHS = {
    "classroom_zip": "classroom.zip",
    "pptx": "classroom.pptx",
    "offline_html": "classroom.html",
    "mp4": "classroom.mp4",
}
CLASSROOM_EXPORT_CONTENT_TYPES = {
    "classroom_zip": "application/zip",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "offline_html": "text/html",
    "mp4": "video/mp4",
}

MAX_CLASSROOM_EXPORT_REPORT_BYTES = 64 * 1024
MAX_EXPORT_BYTES = {
    "classroom_zip": 256 * 1024 * 1024,
    "pptx": 256 * 1024 * 1024,
    "offline_html": 16 * 1024 * 1024,
    "mp4": 2 * 1024 * 1024 * 1024,
}
MAX_TOTAL_EXPORT_BYTES = 3 * 1024 * 1024 * 1024
MAX_ARCHIVE_MEMBERS = 4_096
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 512 * 1024 * 1024
MAX_CLASSROOM_DOCUMENT_BYTES = 16 * 1024 * 1024
MAX_MP4_METADATA_BYTES = 64 * 1024 * 1024

_PUBLIC_ID = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._:-]{0,127}$")
_SHA256 = re.compile(r"^[0-9a-f]{64}$")
_OBSERVED_AT = re.compile(r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}(?:\.\d+)?Z$")
_OOXML_CONTENT_TYPES = "http://schemas.openxmlformats.org/package/2006/content-types"
_PACKAGE_RELATIONSHIPS = "http://schemas.openxmlformats.org/package/2006/relationships"
_OFFICE_RELATIONSHIPS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PRESENTATION = "http://schemas.openxmlformats.org/presentationml/2006/main"


def classroom_exports_command_record() -> dict[str, object]:
    """Return the secret-free logical command recorded in release evidence."""

    return {
        "runner": "python",
        "script": "scripts/classroom_export_probe.py",
        "arguments": ["--profile", "first-release"],
    }


def canonical_classroom_export_report(report: Mapping[str, object]) -> bytes:
    return (
        json.dumps(
            report,
            allow_nan=False,
            ensure_ascii=False,
            separators=(",", ":"),
            sort_keys=True,
        )
        + "\n"
    ).encode("utf-8")


def _exact_json_equal(left: object, right: object) -> bool:
    if type(left) is not type(right):
        return False
    if isinstance(left, dict):
        return set(left) == set(right) and all(
            _exact_json_equal(left[key], right[key]) for key in left
        )
    if isinstance(left, list):
        return len(left) == len(right) and all(
            _exact_json_equal(left_item, right_item)
            for left_item, right_item in zip(left, right, strict=True)
        )
    return left == right


def _valid_public_id(raw: object) -> bool:
    return isinstance(raw, str) and _PUBLIC_ID.fullmatch(raw) is not None


def _valid_sha256(raw: object) -> bool:
    return isinstance(raw, str) and _SHA256.fullmatch(raw) is not None and raw != "0" * 64


def _valid_base_url(raw: object) -> bool:
    if not isinstance(raw, str) or not raw or raw != raw.rstrip("/"):
        return False
    parsed = urlsplit(raw)
    return (
        parsed.scheme in {"http", "https"}
        and bool(parsed.netloc)
        and parsed.username is None
        and parsed.password is None
        and not parsed.query
        and not parsed.fragment
    )


def _valid_observed_at(raw: object) -> bool:
    if not isinstance(raw, str) or _OBSERVED_AT.fullmatch(raw) is None:
        return False
    try:
        datetime.fromisoformat(raw.removesuffix("Z") + "+00:00")
    except ValueError:
        return False
    return True


def _validated_export_relative_path(
    relative_path: object,
    *,
    expected_path: str,
) -> str:
    if not isinstance(relative_path, str) or relative_path != expected_path:
        raise ValueError("classroom export path is invalid")
    posix_path = PurePosixPath(relative_path)
    windows_path = PureWindowsPath(relative_path)
    if (
        posix_path.is_absolute()
        or windows_path.is_absolute()
        or len(posix_path.parts) != 1
        or any(part in {"", ".", ".."} for part in posix_path.parts)
    ):
        raise ValueError("classroom export path is unsafe")
    return relative_path


def _open_export_path(
    artifact_root: Path,
    relative_path: object,
    *,
    expected_path: str,
) -> Path:
    relative_path = _validated_export_relative_path(
        relative_path,
        expected_path=expected_path,
    )

    root = Path(os.path.abspath(os.fspath(artifact_root)))
    reparse_attribute = getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0)
    for component in (*reversed(root.parents), root):
        try:
            component_stat = component.lstat()
        except OSError as exc:
            raise ValueError("classroom export artifact root is invalid") from exc
        if stat.S_ISLNK(component_stat.st_mode) or (
            reparse_attribute
            and getattr(component_stat, "st_file_attributes", 0) & reparse_attribute
        ):
            raise ValueError("classroom export artifact root must not use a symlink or junction")
    try:
        resolved_root = root.resolve(strict=True)
    except (OSError, RuntimeError) as exc:
        raise ValueError("classroom export artifact root is invalid") from exc
    if not resolved_root.is_dir():
        raise ValueError("classroom export artifact root is invalid")
    export_path = resolved_root / relative_path
    try:
        path_stat = export_path.lstat()
    except OSError as exc:
        raise ValueError("classroom export path is invalid") from exc
    if stat.S_ISLNK(path_stat.st_mode):
        raise ValueError("classroom export basename must not be a symlink")
    try:
        resolved_path = export_path.resolve(strict=True)
        resolved_path.relative_to(resolved_root)
        file_stat = resolved_path.stat()
    except (OSError, RuntimeError, ValueError) as exc:
        raise ValueError("classroom export path escapes artifact root") from exc
    if not stat.S_ISREG(file_stat.st_mode):
        raise ValueError("classroom export path is not a regular file")
    return resolved_path


def _hash_file(path: Path, *, max_bytes: int) -> tuple[int, str]:
    file_stat = path.stat()
    if file_stat.st_size <= 0 or file_stat.st_size > max_bytes:
        raise ValueError("classroom export byte length is invalid")
    digest = hashlib.sha256()
    observed = 0
    with path.open("rb") as handle:
        while chunk := handle.read(1024 * 1024):
            observed += len(chunk)
            if observed > max_bytes:
                raise ValueError("classroom export byte length is invalid")
            digest.update(chunk)
    if observed != file_stat.st_size:
        raise ValueError("classroom export byte length changed while reading")
    return observed, digest.hexdigest()


def _hash_handle(handle: BinaryIO, *, max_bytes: int) -> tuple[int, str]:
    try:
        file_stat = os.fstat(handle.fileno())
    except (AttributeError, OSError) as exc:
        raise ValueError("classroom export handle is invalid") from exc
    if (
        not stat.S_ISREG(file_stat.st_mode)
        or file_stat.st_size <= 0
        or file_stat.st_size > max_bytes
    ):
        raise ValueError("classroom export byte length is invalid")
    digest = hashlib.sha256()
    observed = 0
    handle.seek(0)
    try:
        while chunk := handle.read(1024 * 1024):
            observed += len(chunk)
            if observed > max_bytes:
                raise ValueError("classroom export byte length is invalid")
            digest.update(chunk)
    finally:
        handle.seek(0)
    if observed != file_stat.st_size:
        raise ValueError("classroom export byte length changed while reading")
    return observed, digest.hexdigest()


def _validated_archive_members(archive: zipfile.ZipFile) -> dict[str, zipfile.ZipInfo]:
    infos = archive.infolist()
    if len(infos) > MAX_ARCHIVE_MEMBERS:
        raise ValueError("classroom export archive has too many members")
    members: dict[str, zipfile.ZipInfo] = {}
    total_uncompressed = 0
    for info in infos:
        name = info.filename
        path = PurePosixPath(name)
        unix_mode = (info.external_attr >> 16) & 0xFFFF
        if (
            not name
            or "\\" in name
            or path.is_absolute()
            or any(part in {"", ".", ".."} for part in path.parts)
            or stat.S_ISLNK(unix_mode)
            or name in members
        ):
            raise ValueError("classroom export archive contains an unsafe member")
        if info.flag_bits & 0x1:
            raise ValueError("classroom export archive contains an encrypted member")
        total_uncompressed += info.file_size
        if total_uncompressed > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
            raise ValueError("classroom export archive expands too large")
        members[name] = info
    return members


def _safe_archive_members(archive: zipfile.ZipFile) -> dict[str, zipfile.ZipInfo]:
    members = _validated_archive_members(archive)
    if archive.testzip() is not None:
        raise ValueError("classroom export archive failed its CRC check")
    return members


def classroom_export_archive_contains_forbidden_bytes(
    source: BinaryIO,
    *,
    kind: str,
    forbidden: set[bytes],
) -> bool:
    """Scan expanded ZIP/PPTX members without reopening the retained artifact."""

    if kind not in {"classroom_zip", "pptx"} or not forbidden:
        return False
    secrets = {secret for secret in forbidden if isinstance(secret, bytes) and secret}
    if not secrets:
        return False
    overlap = max(len(secret) for secret in secrets) - 1
    expanded_total = 0
    try:
        source.seek(0)
        with zipfile.ZipFile(source) as archive:
            members = _validated_archive_members(archive)
            for info in members.values():
                if info.is_dir():
                    continue
                member_size = 0
                previous = b""
                try:
                    with archive.open(info, "r") as member:
                        while chunk := member.read(1024 * 1024):
                            member_size += len(chunk)
                            expanded_total += len(chunk)
                            if (
                                member_size > info.file_size
                                or expanded_total > MAX_ARCHIVE_UNCOMPRESSED_BYTES
                            ):
                                raise ValueError("classroom export archive expands too large")
                            window = previous + chunk
                            if any(secret in window for secret in secrets):
                                return True
                            previous = window[-overlap:] if overlap > 0 else b""
                except (OSError, RuntimeError, NotImplementedError, zipfile.BadZipFile) as exc:
                    raise ValueError("classroom export archive member is invalid") from exc
                if member_size != info.file_size:
                    raise ValueError("classroom export archive member size changed")
        return False
    except (OSError, zipfile.BadZipFile, zipfile.LargeZipFile) as exc:
        raise ValueError("classroom export archive is invalid") from exc
    finally:
        source.seek(0)


def _read_archive_member(
    archive: zipfile.ZipFile,
    members: Mapping[str, zipfile.ZipInfo],
    name: str,
    *,
    max_bytes: int,
) -> bytes:
    info = members.get(name)
    if info is None or info.is_dir() or info.file_size > max_bytes:
        raise ValueError(f"classroom ZIP {name} is missing or too large")
    try:
        return archive.read(info)
    except (OSError, RuntimeError) as exc:
        raise ValueError(f"classroom ZIP {name} is invalid") from exc


def _validate_classroom_zip(
    source: Path | BinaryIO,
    *,
    classroom_version_id: str,
    document_sha256: str,
) -> None:
    try:
        if not isinstance(source, Path):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            members = _safe_archive_members(archive)
            classroom_body = _read_archive_member(
                archive,
                members,
                "classroom.json",
                max_bytes=MAX_CLASSROOM_DOCUMENT_BYTES,
            )
            if hashlib.sha256(classroom_body).hexdigest() != document_sha256:
                raise ValueError("classroom ZIP classroom.json hash is invalid")
            try:
                classroom = ClassroomDocument.model_validate_json(classroom_body)
            except (ValidationError, ValueError) as exc:
                raise ValueError(
                    "classroom ZIP classroom.json is not a valid ClassroomDocument"
                ) from exc
            if classroom.classroom_version_id != classroom_version_id:
                raise ValueError("classroom ZIP classroom version is invalid")

            media_by_path = {}
            for media in classroom.media_manifest:
                relative_path = media.relative_path
                parsed_path = PurePosixPath(relative_path)
                if (
                    not relative_path
                    or "\\" in relative_path
                    or parsed_path.is_absolute()
                    or any(part in {"", ".", ".."} for part in parsed_path.parts)
                    or relative_path in media_by_path
                ):
                    raise ValueError("classroom ZIP media path is unsafe or duplicated")
                media_by_path[relative_path] = media
            if set(members) != {"classroom.json", *media_by_path}:
                raise ValueError("classroom ZIP members do not match mediaManifest")
            for relative_path, media in media_by_path.items():
                media_body = _read_archive_member(
                    archive,
                    members,
                    relative_path,
                    max_bytes=MAX_ARCHIVE_UNCOMPRESSED_BYTES,
                )
                if len(media_body) != media.size_bytes:
                    raise ValueError("classroom ZIP media size is invalid")
                if hashlib.sha256(media_body).hexdigest() != media.sha256:
                    raise ValueError("classroom ZIP media hash is invalid")
    except (OSError, zipfile.BadZipFile, zipfile.LargeZipFile) as exc:
        raise ValueError("classroom ZIP archive is invalid") from exc
    finally:
        if not isinstance(source, Path):
            source.seek(0)


def _xml_member(
    archive: zipfile.ZipFile,
    members: Mapping[str, zipfile.ZipInfo],
    name: str,
) -> ET.Element:
    info = members.get(name)
    if info is None or info.is_dir() or info.file_size > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
        raise ValueError(f"PPTX {name} is missing")
    try:
        return ET.fromstring(archive.read(info))
    except (OSError, ET.ParseError, RuntimeError) as exc:
        raise ValueError(f"PPTX {name} XML is invalid") from exc


def _relationship_target(base: str, target: str) -> str:
    if not target or "\\" in target or PurePosixPath(target).is_absolute():
        raise ValueError("PPTX relationship target is unsafe")
    resolved = posixpath.normpath(posixpath.join(base, target))
    if resolved == ".." or resolved.startswith("../"):
        raise ValueError("PPTX relationship target is unsafe")
    return resolved


def _validate_pptx(source: Path | BinaryIO, *, document_sha256: str) -> None:
    try:
        if not isinstance(source, Path):
            source.seek(0)
        with zipfile.ZipFile(source) as archive:
            members = _safe_archive_members(archive)
            active_member_markers = (
                "/activex/",
                "/ctrlprops/",
                "/embeddings/",
                "/vbaproject",
            )
            for member_name in members:
                normalized_member = f"/{member_name.casefold()}"
                if any(marker in normalized_member for marker in active_member_markers):
                    raise ValueError("PPTX macro, ActiveX, or OLE active content is invalid")

            content_types = _xml_member(archive, members, "[Content_Types].xml")
            root_rels = _xml_member(archive, members, "_rels/.rels")
            presentation = _xml_member(archive, members, "ppt/presentation.xml")
            presentation_rels = _xml_member(archive, members, "ppt/_rels/presentation.xml.rels")

            if content_types.tag != f"{{{_OOXML_CONTENT_TYPES}}}Types":
                raise ValueError("PPTX content types root is invalid")
            if root_rels.tag != f"{{{_PACKAGE_RELATIONSHIPS}}}Relationships":
                raise ValueError("PPTX package relationships root is invalid")
            if presentation.tag != f"{{{_PRESENTATION}}}presentation":
                raise ValueError("PPTX presentation root is invalid")
            if presentation_rels.tag != f"{{{_PACKAGE_RELATIONSHIPS}}}Relationships":
                raise ValueError("PPTX presentation relationships root is invalid")

            active_type_markers = (
                "macroenabled",
                "vbaproject",
                "activex",
                "oleobject",
            )
            if any(
                any(
                    marker in (node.get("ContentType") or "").casefold()
                    for marker in active_type_markers
                )
                for node in content_types
            ):
                raise ValueError("PPTX macro, ActiveX, or OLE active content is invalid")

            for relationship_path in members:
                if not relationship_path.casefold().endswith(".rels"):
                    continue
                relationships = _xml_member(
                    archive,
                    members,
                    relationship_path,
                )
                if relationships.tag != f"{{{_PACKAGE_RELATIONSHIPS}}}Relationships":
                    raise ValueError("PPTX relationships root is invalid")
                for relationship in relationships.findall(
                    f"{{{_PACKAGE_RELATIONSHIPS}}}Relationship"
                ):
                    if (relationship.get("TargetMode") or "").casefold() == "external":
                        raise ValueError("PPTX external relationship is invalid")
                    relationship_type = (relationship.get("Type") or "").casefold()
                    if any(marker in relationship_type for marker in active_type_markers):
                        raise ValueError("PPTX macro, ActiveX, or OLE active content is invalid")

            overrides = {
                node.get("PartName"): node.get("ContentType")
                for node in content_types.findall(f"{{{_OOXML_CONTENT_TYPES}}}Override")
            }
            if overrides.get("/ppt/presentation.xml") != (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
            ):
                raise ValueError("PPTX presentation content type is invalid")

            office_documents = [
                node
                for node in root_rels.findall(f"{{{_PACKAGE_RELATIONSHIPS}}}Relationship")
                if node.get("Type") == f"{_OFFICE_RELATIONSHIPS}/officeDocument"
                and node.get("TargetMode") != "External"
            ]
            if (
                len(office_documents) != 1
                or _relationship_target("", office_documents[0].get("Target", ""))
                != "ppt/presentation.xml"
            ):
                raise ValueError("PPTX presentation relationship is invalid")

            relationship_by_id: dict[str, ET.Element] = {}
            for node in presentation_rels.findall(f"{{{_PACKAGE_RELATIONSHIPS}}}Relationship"):
                relationship_id = node.get("Id")
                if not relationship_id or relationship_id in relationship_by_id:
                    raise ValueError("PPTX presentation relationship is invalid")
                relationship_by_id[relationship_id] = node

            slide_ids = presentation.findall(f".//{{{_PRESENTATION}}}sldId")
            if not slide_ids:
                raise ValueError("PPTX must contain at least one slide")
            slide_paths: list[str] = []
            seen_slide_ids: set[int] = set()
            seen_slide_relationships: set[str] = set()
            for slide_id in slide_ids:
                raw_slide_id = slide_id.get("id")
                relationship_id = slide_id.get(f"{{{_OFFICE_RELATIONSHIPS}}}id")
                try:
                    numeric_slide_id = int(raw_slide_id or "")
                except ValueError as exc:
                    raise ValueError("PPTX slide identity is invalid") from exc
                if (
                    numeric_slide_id < 256
                    or numeric_slide_id in seen_slide_ids
                    or not relationship_id
                    or relationship_id in seen_slide_relationships
                ):
                    raise ValueError("PPTX slide identity is invalid")
                seen_slide_ids.add(numeric_slide_id)
                seen_slide_relationships.add(relationship_id)
                relationship = relationship_by_id.get(relationship_id or "")
                if (
                    relationship is None
                    or relationship.get("Type") != f"{_OFFICE_RELATIONSHIPS}/slide"
                    or relationship.get("TargetMode") == "External"
                ):
                    raise ValueError("PPTX slide relationship is invalid")
                slide_path = _relationship_target("ppt", relationship.get("Target", ""))
                if (
                    not slide_path.startswith("ppt/slides/")
                    or slide_path not in members
                    or overrides.get(f"/{slide_path}")
                    != "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
                ):
                    raise ValueError("PPTX slide part is invalid")
                slide_paths.append(slide_path)

            for slide_path in slide_paths:
                slide = _xml_member(archive, members, slide_path)
                if slide.tag != f"{{{_PRESENTATION}}}sld":
                    raise ValueError("PPTX slide XML is invalid")
    except (OSError, zipfile.BadZipFile, zipfile.LargeZipFile) as exc:
        raise ValueError("PPTX archive is invalid") from exc
    try:
        if not isinstance(source, Path):
            source.seek(0)
        opened = Presentation(str(source) if isinstance(source, Path) else source)
    except (OSError, KeyError, ValueError, PythonPptxError, zipfile.BadZipFile) as exc:
        raise ValueError("PPTX could not be opened by python-pptx") from exc
    finally:
        if not isinstance(source, Path):
            source.seek(0)
    if len(opened.slides) < 1:
        raise ValueError("PPTX must contain at least one slide")
    slide_width = int(opened.slide_width or 0)
    slide_height = int(opened.slide_height or 0)
    if slide_width <= 0 or slide_height <= 0:
        raise ValueError("PPTX slide bounds are invalid")
    hash_is_visible = False
    for slide in opened.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if not isinstance(text, str) or document_sha256 not in text:
                continue
            left = int(shape.left or 0)
            top = int(shape.top or 0)
            width = int(shape.width or 0)
            height = int(shape.height or 0)
            if (
                left >= 0
                and top >= 0
                and width > 0
                and height > 0
                and left + width <= slide_width
                and top + height <= slide_height
            ):
                hash_is_visible = True
                break
        if hash_is_visible:
            break
    if not hash_is_visible:
        raise ValueError("PPTX visible in-bounds shape text does not bind the document hash")


_HTML_RESOURCE_ATTRIBUTES = {"action", "href", "src", "srcset"}
_HTML_ACTIVE_ELEMENTS = {
    "base",
    "embed",
    "form",
    "frame",
    "frameset",
    "iframe",
    "link",
    "object",
}
_CSS_RESOURCE_REFERENCE = re.compile(r"(?:@import\b|url\s*\()", re.IGNORECASE)


class _OfflineHTMLInspector(HTMLParser):
    _VOID_ELEMENTS = {
        "area",
        "base",
        "br",
        "col",
        "embed",
        "hr",
        "img",
        "input",
        "link",
        "meta",
        "param",
        "source",
        "track",
        "wbr",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.stack: list[str] = []
        self.doctype_count = 0
        self.html_count = 0
        self.head_count = 0
        self.body_count = 0
        self.main_count = 0
        self.head_closed = False
        self.body_closed = False
        self.body_document_sha256: str | None = None
        self.classroom_script_count = 0
        self.classroom_script_parts: list[str] = []
        self.collecting_classroom_script = False
        self.collecting_style = False
        self.style_parts: list[str] = []
        self.non_whitespace_outside_document = False

    def handle_decl(self, decl: str) -> None:
        if decl.casefold() != "doctype html" or self.stack or self.html_count or self.doctype_count:
            raise ValueError("offline HTML5 doctype is invalid")
        self.doctype_count = 1

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        lowered = tag.lower()
        names = [name.lower() for name, _value in attrs]
        if len(names) != len(set(names)):
            raise ValueError("offline HTML attributes are invalid")
        attributes = {name.lower(): value for name, value in attrs}
        if any(name.startswith("on") for name in attributes):
            raise ValueError("offline HTML active event attribute is invalid")
        if (
            lowered == "meta"
            and (attributes.get("http-equiv") or "").strip().casefold() == "refresh"
        ):
            raise ValueError("offline HTML active meta refresh is invalid")
        is_classroom_script = lowered == "script" and attributes == {
            "id": "classroom",
            "type": "application/json",
        }
        if lowered == "script" and not is_classroom_script:
            raise ValueError("offline HTML active non-classroom script is invalid")
        if lowered in _HTML_ACTIVE_ELEMENTS:
            raise ValueError(f"offline HTML active {lowered} element is invalid")
        resource_attributes = _HTML_RESOURCE_ATTRIBUTES.intersection(attributes)
        if resource_attributes:
            attribute = sorted(resource_attributes)[0]
            raise ValueError(f"offline HTML {attribute} resource attribute is invalid")
        if any(
            isinstance(value, str) and value.lstrip().startswith("//")
            for value in attributes.values()
        ):
            raise ValueError("offline HTML protocol-relative external URL is invalid")
        style_attribute = attributes.get("style")
        if isinstance(style_attribute, str) and _CSS_RESOURCE_REFERENCE.search(style_attribute):
            raise ValueError("offline HTML CSS url/import resource is invalid")

        if not self.stack and lowered != "html":
            raise ValueError("offline HTML document structure is invalid")
        if lowered == "html":
            self.html_count += 1
            if self.doctype_count != 1:
                raise ValueError("offline HTML5 doctype is missing or invalid")
            if self.stack or self.html_count != 1:
                raise ValueError("offline HTML root is invalid")
        elif lowered == "head":
            self.head_count += 1
            if self.stack != ["html"] or self.head_count != 1 or self.body_count:
                raise ValueError("offline HTML head structure is invalid")
        elif lowered == "body":
            self.body_count += 1
            if (
                self.stack != ["html"]
                or self.head_count != 1
                or not self.head_closed
                or self.body_count != 1
            ):
                raise ValueError("offline HTML body structure is invalid")
            self.body_document_sha256 = attributes.get("data-document-sha256")
        elif lowered == "main":
            self.main_count += 1
            if self.stack != ["html", "body"] or self.main_count != 1:
                raise ValueError("offline HTML main structure is invalid")
        elif lowered == "style":
            if "head" not in self.stack:
                raise ValueError("offline HTML style must be inside head")
            self.collecting_style = True
            self.style_parts = []
        elif is_classroom_script:
            if self.stack != ["html", "body"]:
                raise ValueError("offline HTML classroom script is outside body")
            self.classroom_script_count += 1
            self.collecting_classroom_script = True
        if self.stack == ["html"] and lowered not in {"body", "head"}:
            raise ValueError("offline HTML document structure is invalid")
        if lowered not in self._VOID_ELEMENTS:
            self.stack.append(lowered)

    def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        self.handle_starttag(tag, attrs)
        if tag.lower() not in self._VOID_ELEMENTS:
            self.handle_endtag(tag)

    def handle_endtag(self, tag: str) -> None:
        lowered = tag.lower()
        if lowered in self._VOID_ELEMENTS or not self.stack or self.stack[-1] != lowered:
            raise ValueError("offline HTML tags are not properly closed")
        self.stack.pop()
        if lowered == "head":
            self.head_closed = True
        elif lowered == "body":
            self.body_closed = True
        if lowered == "script" and self.collecting_classroom_script:
            self.collecting_classroom_script = False
        if lowered == "style" and self.collecting_style:
            self.collecting_style = False
            if _CSS_RESOURCE_REFERENCE.search("".join(self.style_parts)):
                raise ValueError("offline HTML CSS url/import resource is invalid")

    def handle_data(self, data: str) -> None:
        if self.collecting_classroom_script:
            self.classroom_script_parts.append(data)
        elif self.collecting_style:
            self.style_parts.append(data)
        elif not self.stack and data.strip():
            self.non_whitespace_outside_document = True


def _validate_offline_html(
    source: Path | BinaryIO,
    *,
    classroom_version_id: str,
    document_sha256: str,
) -> None:
    try:
        if isinstance(source, Path):
            body = source.read_text(encoding="utf-8")
        else:
            source.seek(0)
            try:
                body = source.read().decode("utf-8", errors="strict")
            finally:
                source.seek(0)
    except (OSError, UnicodeError) as exc:
        raise ValueError("offline HTML must be valid UTF-8") from exc
    inspector = _OfflineHTMLInspector()
    try:
        inspector.feed(body)
        inspector.close()
    except ValueError as exc:
        raise ValueError(f"offline HTML is invalid: {exc}") from exc
    except RuntimeError as exc:
        raise ValueError("offline HTML is invalid") from exc
    if (
        inspector.stack
        or inspector.non_whitespace_outside_document
        or inspector.doctype_count != 1
        or inspector.html_count != 1
        or inspector.head_count != 1
        or inspector.body_count != 1
        or not inspector.head_closed
        or not inspector.body_closed
    ):
        raise ValueError("offline HTML document structure is invalid")
    if inspector.main_count != 1:
        raise ValueError("offline HTML main element is missing or duplicated")
    if (
        inspector.body_document_sha256 is not None
        and inspector.body_document_sha256 != document_sha256
    ):
        raise ValueError("offline HTML document hash binding is invalid")
    if inspector.classroom_script_count != 1:
        raise ValueError("offline HTML classroom script is missing or duplicated")
    classroom_body = "".join(inspector.classroom_script_parts)
    try:
        classroom = ClassroomDocument.model_validate_json(classroom_body)
    except (ValidationError, ValueError) as exc:
        raise ValueError("offline HTML classroom script is not a valid ClassroomDocument") from exc
    if classroom.classroom_version_id != classroom_version_id:
        raise ValueError("offline HTML classroom version is invalid")
    if hashlib.sha256(canonical_json_bytes(classroom)).hexdigest() != document_sha256:
        raise ValueError("offline HTML classroom document hash is invalid")


def _parse_mp4_box_payload(payload: bytes, *, context: str) -> list[tuple[bytes, bytes]]:
    offset = 0
    boxes: list[tuple[bytes, bytes]] = []
    while offset < len(payload):
        if len(payload) - offset < 8:
            raise ValueError(f"MP4 {context} child box header is truncated")
        size_32, box_type = struct.unpack_from(">I4s", payload, offset)
        header_size = 8
        if size_32 == 0:
            raise ValueError(f"MP4 {context} child box size 0 is invalid")
        if size_32 == 1:
            if len(payload) - offset < 16:
                raise ValueError(f"MP4 {context} extended child size is truncated")
            box_size = struct.unpack_from(">Q", payload, offset + 8)[0]
            header_size = 16
        else:
            box_size = size_32
        if box_size < header_size or box_size > len(payload) - offset:
            raise ValueError(f"MP4 {context} child box is truncated or invalid")
        if any(character < 0x20 or character > 0x7E for character in box_type):
            raise ValueError(f"MP4 {context} child box type is invalid")
        start = offset + header_size
        end = offset + box_size
        boxes.append((box_type, payload[start:end]))
        offset = end
    return boxes


def _mp4_full_box_version(
    payload: bytes,
    *,
    context: str,
    minimum_v0: int,
    minimum_v1: int,
) -> int:
    if not payload or payload[0] not in {0, 1}:
        raise ValueError(f"MP4 {context} version is invalid")
    version = payload[0]
    minimum_bytes = minimum_v0 if version == 0 else minimum_v1
    if len(payload) < minimum_bytes:
        raise ValueError(f"MP4 {context} payload is too short")
    return version


def _validate_mp4_timing(
    payload: bytes,
    *,
    context: str,
    minimum_v0: int,
    minimum_v1: int,
) -> None:
    version = _mp4_full_box_version(
        payload,
        context=context,
        minimum_v0=minimum_v0,
        minimum_v1=minimum_v1,
    )
    if version == 0:
        timescale = struct.unpack_from(">I", payload, 12)[0]
        duration = struct.unpack_from(">I", payload, 16)[0]
    else:
        timescale = struct.unpack_from(">I", payload, 20)[0]
        duration = struct.unpack_from(">Q", payload, 24)[0]
    if timescale == 0:
        raise ValueError(f"MP4 {context} timescale must be nonzero")
    if duration == 0:
        raise ValueError(f"MP4 {context} duration must be nonzero")


def _validate_tkhd(payload: bytes) -> int:
    version = _mp4_full_box_version(
        payload,
        context="tkhd",
        minimum_v0=84,
        minimum_v1=96,
    )
    track_id_offset = 12 if version == 0 else 20
    track_id = struct.unpack_from(">I", payload, track_id_offset)[0]
    if track_id == 0:
        raise ValueError("MP4 tkhd track ID must be nonzero")
    return track_id


def _mp4_counted_table_entries(
    payload: bytes,
    *,
    context: str,
    entry_bytes: int,
) -> int:
    if len(payload) < 8 or payload[:4] != b"\x00\x00\x00\x00":
        raise ValueError(f"MP4 {context} FullBox header is invalid")
    entry_count = struct.unpack_from(">I", payload, 4)[0]
    if entry_count == 0:
        raise ValueError(f"MP4 {context} entry table is empty")
    if len(payload) != 8 + entry_count * entry_bytes:
        raise ValueError(f"MP4 {context} entry table length is invalid")
    return entry_count


def _validate_video_sample_table(stbl_boxes: list[tuple[bytes, bytes]]) -> None:
    def one_payload(kind: bytes) -> bytes:
        matches = [payload for box_type, payload in stbl_boxes if box_type == kind]
        if len(matches) != 1:
            raise ValueError(f"MP4 video stbl must contain exactly one {kind.decode()} box")
        return matches[0]

    stsd = one_payload(b"stsd")
    if len(stsd) < 8 or stsd[:4] != b"\x00\x00\x00\x00":
        raise ValueError("MP4 stsd FullBox header is invalid")
    stsd_entry_count = struct.unpack_from(">I", stsd, 4)[0]
    if stsd_entry_count == 0:
        raise ValueError("MP4 stsd sample entry table is empty")
    sample_entries = _parse_mp4_box_payload(stsd[8:], context="stsd")
    if len(sample_entries) != stsd_entry_count:
        raise ValueError("MP4 stsd sample entry count is invalid")
    supported_visual_entries = {
        b"av01",
        b"avc1",
        b"avc3",
        b"hev1",
        b"hvc1",
        b"vp08",
        b"vp09",
    }
    valid_visual_entry = False
    for sample_type, sample_payload in sample_entries:
        if sample_type not in supported_visual_entries:
            continue
        if len(sample_payload) < 78:
            raise ValueError("MP4 stsd visual sample entry is too short")
        data_reference_index = struct.unpack_from(">H", sample_payload, 6)[0]
        width, height = struct.unpack_from(">HH", sample_payload, 24)
        if data_reference_index == 0 or width == 0 or height == 0:
            raise ValueError("MP4 stsd visual sample entry is invalid")
        valid_visual_entry = True
    if not valid_visual_entry:
        raise ValueError("MP4 stsd lacks a supported visual sample entry")

    stts = one_payload(b"stts")
    stts_count = _mp4_counted_table_entries(
        stts,
        context="stts",
        entry_bytes=8,
    )
    for index in range(stts_count):
        sample_count, sample_delta = struct.unpack_from(">II", stts, 8 + index * 8)
        if sample_count == 0 or sample_delta == 0:
            raise ValueError("MP4 stts entry is invalid")

    stsc = one_payload(b"stsc")
    stsc_count = _mp4_counted_table_entries(
        stsc,
        context="stsc",
        entry_bytes=12,
    )
    for index in range(stsc_count):
        entry = struct.unpack_from(">III", stsc, 8 + index * 12)
        if any(value == 0 for value in entry):
            raise ValueError("MP4 stsc entry is invalid")

    stsz = one_payload(b"stsz")
    if len(stsz) < 12 or stsz[:4] != b"\x00\x00\x00\x00":
        raise ValueError("MP4 stsz FullBox header is invalid")
    sample_size, sample_count = struct.unpack_from(">II", stsz, 4)
    if sample_count == 0:
        raise ValueError("MP4 stsz sample table is empty")
    expected_stsz_bytes = 12 if sample_size else 12 + sample_count * 4
    if len(stsz) != expected_stsz_bytes:
        raise ValueError("MP4 stsz sample table length is invalid")
    if sample_size == 0 and any(
        struct.unpack_from(">I", stsz, 12 + index * 4)[0] == 0 for index in range(sample_count)
    ):
        raise ValueError("MP4 stsz sample size is invalid")

    chunk_offsets = [
        (box_type, payload) for box_type, payload in stbl_boxes if box_type in {b"stco", b"co64"}
    ]
    if len(chunk_offsets) != 1:
        raise ValueError("MP4 video stbl must contain exactly one stco or co64 box")
    chunk_type, chunk_payload = chunk_offsets[0]
    entry_bytes = 4 if chunk_type == b"stco" else 8
    chunk_count = _mp4_counted_table_entries(
        chunk_payload,
        context=chunk_type.decode(),
        entry_bytes=entry_bytes,
    )
    chunk_format = ">I" if entry_bytes == 4 else ">Q"
    if any(
        struct.unpack_from(chunk_format, chunk_payload, 8 + index * entry_bytes)[0] == 0
        for index in range(chunk_count)
    ):
        raise ValueError("MP4 chunk offset entry is invalid")


def _validate_mp4(source: Path | BinaryIO) -> None:
    file_size = (
        source.stat().st_size if isinstance(source, Path) else os.fstat(source.fileno()).st_size
    )
    offset = 0
    box_types: list[bytes] = []
    moov_payloads: list[bytes] = []
    mdat_count = 0
    try:
        if isinstance(source, Path):
            handle_context = source.open("rb")
        else:
            source.seek(0)
            handle_context = source
        handle = handle_context
        try:
            while offset < file_size:
                header = handle.read(8)
                if len(header) != 8:
                    raise ValueError("MP4 top-level box header is truncated")
                size_32, box_type = struct.unpack(">I4s", header)
                header_size = 8
                if size_32 == 0:
                    raise ValueError("MP4 top-level box size 0 is not allowed")
                if size_32 == 1:
                    extended = handle.read(8)
                    if len(extended) != 8:
                        raise ValueError("MP4 extended size is truncated")
                    box_size = struct.unpack(">Q", extended)[0]
                    header_size = 16
                    if box_size < header_size:
                        raise ValueError("MP4 extended box size is invalid")
                else:
                    box_size = size_32
                    if box_size < header_size:
                        raise ValueError("MP4 top-level box size is invalid")
                if any(character < 0x20 or character > 0x7E for character in box_type):
                    raise ValueError("MP4 top-level box type is invalid")
                if box_size > file_size - offset:
                    raise ValueError("MP4 top-level box is truncated")
                payload_size = box_size - header_size
                if box_type == b"ftyp":
                    if payload_size < 8 or (payload_size - 8) % 4:
                        raise ValueError("MP4 ftyp box is invalid")
                    payload = handle.read(payload_size)
                    if len(payload) != payload_size or any(
                        character < 0x20 or character > 0x7E
                        for character in payload[:4] + payload[8:]
                    ):
                        raise ValueError("MP4 ftyp box is invalid")
                elif box_type == b"moov":
                    if payload_size <= 0 or payload_size > MAX_MP4_METADATA_BYTES:
                        raise ValueError("MP4 moov payload is empty or too large")
                    payload = handle.read(payload_size)
                    if len(payload) != payload_size:
                        raise ValueError("MP4 moov payload is truncated")
                    moov_payloads.append(payload)
                elif box_type == b"mdat":
                    if payload_size <= 0:
                        raise ValueError("MP4 mdat payload is empty")
                    handle.seek(payload_size, os.SEEK_CUR)
                    mdat_count += 1
                else:
                    handle.seek(payload_size, os.SEEK_CUR)
                offset += box_size
                box_types.append(box_type)
                if len(box_types) > 100_000:
                    raise ValueError("MP4 contains too many top-level boxes")
        finally:
            if isinstance(source, Path):
                handle.close()
            else:
                source.seek(0)
    except OSError as exc:
        raise ValueError("MP4 could not be opened") from exc
    if offset != file_size:
        raise ValueError("MP4 top-level boxes are incomplete")
    if not box_types or box_types[0] != b"ftyp" or box_types.count(b"ftyp") != 1:
        raise ValueError("MP4 ftyp box is missing or misplaced")
    if box_types.count(b"moov") != 1 or len(moov_payloads) != 1:
        raise ValueError("MP4 must contain exactly one non-empty moov box")
    if b"mdat" not in box_types or mdat_count < 1:
        raise ValueError("MP4 non-empty mdat box is missing")
    moov_boxes = _parse_mp4_box_payload(moov_payloads[0], context="moov")
    mvhd_payloads = [payload for box_type, payload in moov_boxes if box_type == b"mvhd"]
    if len(mvhd_payloads) != 1:
        raise ValueError("MP4 moov must contain exactly one mvhd box")
    _validate_mp4_timing(
        mvhd_payloads[0],
        context="mvhd",
        minimum_v0=100,
        minimum_v1=112,
    )
    trak_payloads = [payload for box_type, payload in moov_boxes if box_type == b"trak"]
    if not trak_payloads:
        raise ValueError("MP4 moov must contain at least one trak box")
    video_track_found = False
    seen_track_ids: set[int] = set()
    for trak_payload in trak_payloads:
        if not trak_payload:
            raise ValueError("MP4 trak payload is empty")
        trak_boxes = _parse_mp4_box_payload(trak_payload, context="trak")
        tkhd_payloads = [payload for box_type, payload in trak_boxes if box_type == b"tkhd"]
        if len(tkhd_payloads) != 1:
            raise ValueError("MP4 trak must contain one non-empty tkhd box")
        track_id = _validate_tkhd(tkhd_payloads[0])
        if track_id in seen_track_ids:
            raise ValueError("MP4 tkhd track ID is duplicated")
        seen_track_ids.add(track_id)
        mdia_payloads = [payload for box_type, payload in trak_boxes if box_type == b"mdia"]
        if len(mdia_payloads) != 1 or not mdia_payloads[0]:
            raise ValueError("MP4 trak must contain one non-empty mdia box")
        mdia_boxes = _parse_mp4_box_payload(mdia_payloads[0], context="mdia")
        mdhd_payloads = [payload for box_type, payload in mdia_boxes if box_type == b"mdhd"]
        if len(mdhd_payloads) != 1:
            raise ValueError("MP4 mdia must contain one valid mdhd box")
        _validate_mp4_timing(
            mdhd_payloads[0],
            context="mdhd",
            minimum_v0=24,
            minimum_v1=36,
        )
        hdlr_payloads = [payload for box_type, payload in mdia_boxes if box_type == b"hdlr"]
        if len(hdlr_payloads) != 1 or len(hdlr_payloads[0]) < 24:
            raise ValueError("MP4 mdia must contain one valid hdlr box")
        hdlr_payload = hdlr_payloads[0]
        if hdlr_payload[0] != 0:
            raise ValueError("MP4 hdlr version is invalid")
        minf_payloads = [payload for box_type, payload in mdia_boxes if box_type == b"minf"]
        if len(minf_payloads) != 1 or not minf_payloads[0]:
            raise ValueError("MP4 mdia must contain one non-empty minf box")
        minf_boxes = _parse_mp4_box_payload(minf_payloads[0], context="minf")
        stbl_payloads = [payload for box_type, payload in minf_boxes if box_type == b"stbl"]
        if len(stbl_payloads) != 1 or not stbl_payloads[0]:
            raise ValueError("MP4 minf must contain one non-empty stbl box")
        stbl_boxes = _parse_mp4_box_payload(stbl_payloads[0], context="stbl")
        if not stbl_boxes:
            raise ValueError("MP4 stbl must contain at least one box")
        if hdlr_payload[8:12] == b"vide":
            _validate_video_sample_table(stbl_boxes)
            video_track_found = True
    if not video_track_found:
        raise ValueError("MP4 must contain at least one video track")


def _parse_export_metadata(
    exports: object,
    *,
    artifact_root: Path,
    artifact_handles: Mapping[str, BinaryIO] | None,
    classroom_version_id: str,
    document_sha256: str,
) -> dict[str, object]:
    if not isinstance(exports, dict) or set(exports) != set(CLASSROOM_EXPORT_KINDS):
        raise ValueError("classroom report exports are invalid")
    seen_export_ids: set[str] = set()
    seen_job_ids: set[str] = set()
    if artifact_handles is not None and set(artifact_handles) != set(CLASSROOM_EXPORT_KINDS):
        raise ValueError("classroom export artifact handles are invalid")
    seen_paths: set[str] = set()
    total_bytes = 0
    for kind in CLASSROOM_EXPORT_KINDS:
        export = exports.get(kind)
        if not isinstance(export, dict) or set(export) != {
            "exportId",
            "jobId",
            "status",
            "progressPercent",
            "downloadReady",
            "relativePath",
            "contentType",
            "contentDisposition",
            "byteLength",
            "sha256",
        }:
            raise ValueError("classroom export metadata is invalid")
        export_id = export.get("exportId")
        job_id = export.get("jobId")
        if (
            not _valid_public_id(export_id)
            or not _valid_public_id(job_id)
            or export_id in seen_export_ids
            or job_id in seen_job_ids
        ):
            raise ValueError("classroom export identity is invalid")
        seen_export_ids.add(export_id)
        seen_job_ids.add(job_id)
        if export.get("status") != "succeeded":
            raise ValueError("classroom export status is invalid")
        if type(export.get("progressPercent")) is not int or export["progressPercent"] != 100:
            raise ValueError("classroom export progress is invalid")
        if export.get("downloadReady") is not True:
            raise ValueError("classroom export download readiness is invalid")
        expected_path = CLASSROOM_EXPORT_PATHS[kind]
        if export.get("contentType") != CLASSROOM_EXPORT_CONTENT_TYPES[kind]:
            raise ValueError("classroom export content type is invalid")
        encoded_basename = quote(expected_path, safe="")
        if export.get("contentDisposition") != (f"attachment; filename*=UTF-8''{encoded_basename}"):
            raise ValueError("classroom export content disposition is invalid")
        declared_size = export.get("byteLength")
        if type(declared_size) is not int or declared_size <= 0:
            raise ValueError("classroom export byte length is invalid")
        if not _valid_sha256(export.get("sha256")):
            raise ValueError("classroom export SHA-256 is invalid")
        relative_path = _validated_export_relative_path(
            export.get("relativePath"),
            expected_path=expected_path,
        )
        if relative_path in seen_paths:
            raise ValueError("classroom export paths must be distinct")
        seen_paths.add(relative_path)
        if artifact_handles is None:
            source: Path | BinaryIO = _open_export_path(
                Path(artifact_root),
                relative_path,
                expected_path=expected_path,
            )
            observed_size, observed_sha256 = _hash_file(
                source,
                max_bytes=MAX_EXPORT_BYTES[kind],
            )
        else:
            source = artifact_handles[kind]
            observed_size, observed_sha256 = _hash_handle(
                source,
                max_bytes=MAX_EXPORT_BYTES[kind],
            )
        if observed_size != declared_size:
            raise ValueError("classroom export byte length does not match the artifact")
        if observed_sha256 != export.get("sha256"):
            raise ValueError("classroom export SHA-256 does not match the artifact")
        total_bytes += observed_size
        if total_bytes > MAX_TOTAL_EXPORT_BYTES:
            raise ValueError("classroom export total is too large")

        if kind == "classroom_zip":
            _validate_classroom_zip(
                source,
                classroom_version_id=classroom_version_id,
                document_sha256=document_sha256,
            )
        elif kind == "pptx":
            _validate_pptx(source, document_sha256=document_sha256)
        elif kind == "offline_html":
            _validate_offline_html(
                source,
                classroom_version_id=classroom_version_id,
                document_sha256=document_sha256,
            )
        else:
            _validate_mp4(source)
    return exports


def parse_classroom_export_report(
    body: bytes,
    *,
    artifact_root: Path,
    artifact_handles: Mapping[str, BinaryIO] | None = None,
    candidate: Mapping[str, object],
    release_run: Mapping[str, str],
    expected_base_url: str,
) -> dict[str, object]:
    """Parse a canonical report against one fixed artifact snapshot.

    Writers and verifiers pass retained file handles. Offline callers may omit
    ``artifact_handles`` and validate the fixed paths under ``artifact_root``.
    """

    if not isinstance(body, bytes) or not body or len(body) > MAX_CLASSROOM_EXPORT_REPORT_BYTES:
        raise ValueError("classroom export report size is invalid")
    try:
        report = json.loads(body)
    except (UnicodeError, json.JSONDecodeError) as exc:
        raise ValueError("classroom export report is invalid") from exc
    if not isinstance(report, dict) or set(report) != {
        "schemaVersion",
        "producer",
        "candidate",
        "releaseRun",
        "observedAt",
        "baseUrl",
        "tenantId",
        "classroomVersionId",
        "documentSha256",
        "exports",
    }:
        raise ValueError("classroom export report is invalid")
    try:
        canonical = canonical_classroom_export_report(report)
    except (TypeError, ValueError) as exc:
        raise ValueError("classroom export report is invalid") from exc
    if canonical != body:
        raise ValueError("classroom export report is not canonical")
    if (
        type(report.get("schemaVersion")) is not int
        or report["schemaVersion"] != CLASSROOM_EXPORT_SCHEMA_VERSION
    ):
        raise ValueError("classroom export report schema is invalid")
    if report.get("producer") != CLASSROOM_EXPORT_PRODUCER:
        raise ValueError("classroom export report producer is invalid")
    if not _exact_json_equal(report.get("candidate"), dict(candidate)) or not (
        _exact_json_equal(report.get("releaseRun"), dict(release_run))
    ):
        raise ValueError("classroom export report release binding is invalid")
    if not _valid_observed_at(report.get("observedAt")):
        raise ValueError("classroom export report timestamp is invalid")
    if not _valid_base_url(report.get("baseUrl")) or report.get("baseUrl") != expected_base_url:
        raise ValueError("classroom export report URL is invalid")
    tenant_id = report.get("tenantId")
    if not _valid_public_id(tenant_id):
        raise ValueError("classroom export tenant is invalid")
    classroom_version_id = report.get("classroomVersionId")
    if not _valid_public_id(classroom_version_id):
        raise ValueError("classroom export classroom version is invalid")
    document_sha256 = report.get("documentSha256")
    if not _valid_sha256(document_sha256):
        raise ValueError("classroom export document hash is invalid")
    _parse_export_metadata(
        report.get("exports"),
        artifact_root=Path(artifact_root),
        artifact_handles=artifact_handles,
        classroom_version_id=classroom_version_id,
        document_sha256=document_sha256,
    )
    return report


def derive_classroom_export_checks(
    body: bytes,
    *,
    artifact_root: Path,
    artifact_handles: Mapping[str, BinaryIO] | None = None,
    candidate: Mapping[str, object],
    release_run: Mapping[str, str],
    expected_base_url: str,
) -> dict[str, bool]:
    """Reparse one fixed artifact snapshot before deriving the checks."""

    parse_classroom_export_report(
        body,
        artifact_root=artifact_root,
        artifact_handles=artifact_handles,
        candidate=candidate,
        release_run=release_run,
        expected_base_url=expected_base_url,
    )
    return {
        "zipOpened": True,
        "pptxOpened": True,
        "offlineHtmlOpened": True,
        "mp4Opened": True,
    }
