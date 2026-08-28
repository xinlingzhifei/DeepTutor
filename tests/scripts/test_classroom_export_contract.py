from __future__ import annotations

from copy import deepcopy
from functools import cache
import hashlib
import importlib.util
import io
import json
from pathlib import Path
import struct
import sys
import xml.etree.ElementTree as ET
import zipfile

from pptx import Presentation
import pytest

from deeptutor.teaching.contracts import ClassroomDocument, canonical_json_bytes
from tests.teaching.test_contracts import valid_classroom_document

ROOT = Path(__file__).resolve().parents[2]
CLASSROOM_VERSION_ID = "classroom-version-immutable-001"
BASE_URL = "https://classroom.example.test"
OBSERVED_AT = "2026-08-28T12:34:56.123Z"
EXPORT_PATHS = {
    "classroom_zip": "classroom.zip",
    "pptx": "classroom.pptx",
    "offline_html": "classroom.html",
    "mp4": "classroom.mp4",
}
CONTENT_TYPES = {
    "classroom_zip": "application/zip",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "offline_html": "text/html",
    "mp4": "video/mp4",
}
CLASSROOM_MEDIA = {"media/voice.mp3": b"first-release-audio"}


def _valid_classroom_body(
    *,
    classroom_version_id: str = CLASSROOM_VERSION_ID,
    media_files: dict[str, bytes] | None = None,
) -> bytes:
    if media_files is None:
        media_files = CLASSROOM_MEDIA
    payload = valid_classroom_document()
    payload["classroom_version_id"] = classroom_version_id
    payload["media_manifest"] = [
        {
            "media_id": "media-1",
            "relative_path": relative_path,
            "mime_type": "audio/mpeg",
            "sha256": hashlib.sha256(body).hexdigest(),
            "size_bytes": len(body),
        }
        for relative_path, body in media_files.items()
    ]
    return canonical_json_bytes(ClassroomDocument.model_validate(payload))


CLASSROOM_BODY = _valid_classroom_body()
DOCUMENT_SHA256 = hashlib.sha256(CLASSROOM_BODY).hexdigest()


@cache
def _module():
    path = ROOT / "scripts" / "classroom_export_contract.py"
    assert path.is_file(), "classroom export contract must exist"
    spec = importlib.util.spec_from_file_location("classroom_export_contract_under_test", path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def _candidate() -> dict[str, object]:
    return {
        "sourceRepository": "xinlingzhifei/DeepTutor",
        "sourceHead": "a" * 40,
        "releaseTag": "yfeistai-first-release-20260828-aaaaaaaa",
        "openmaicHead": "0cf2a330411681190e89f48e20f305345ff99f87",
        "imageDigests": {
            "deeptutor": "sha256:" + "1" * 64,
            "openmaic": "sha256:" + "2" * 64,
            "openmaic_render": "sha256:" + "3" * 64,
        },
    }


def _release_run() -> dict[str, str]:
    return {"runId": "run-classroom-export", "environmentId": "environment-export"}


def _canonical_json(value: object) -> bytes:
    return (
        json.dumps(
            value,
            allow_nan=False,
            ensure_ascii=False,
            separators=(",", ":"),
            sort_keys=True,
        )
        + "\n"
    ).encode()


def _zip_bytes(files: dict[str, bytes]) -> bytes:
    stream = io.BytesIO()
    with zipfile.ZipFile(stream, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, body in files.items():
            archive.writestr(name, body)
    return stream.getvalue()


def _classroom_zip_bytes(
    *,
    classroom_body: bytes | None = None,
    media_files: dict[str, bytes] | None = None,
    extra_files: dict[str, bytes] | None = None,
) -> bytes:
    if classroom_body is None:
        classroom_body = CLASSROOM_BODY
    if media_files is None:
        media_files = CLASSROOM_MEDIA
    files = {"classroom.json": classroom_body, **media_files}
    if extra_files:
        files.update(extra_files)
    return _zip_bytes(files)


def _pptx_bytes(*, document_sha256: str = DOCUMENT_SHA256) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "Verified classroom export"
    slide.placeholders[1].text = f"Document SHA-256: {document_sha256}"
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _replace_zip_member(body: bytes, name: str, replacement: bytes) -> bytes:
    files: dict[str, bytes] = {}
    with zipfile.ZipFile(io.BytesIO(body)) as source:
        for info in source.infolist():
            files[info.filename] = replacement if info.filename == name else source.read(info)
    return _zip_bytes(files)


def _pptx_with_hash_only_in_xml_attribute() -> bytes:
    body = _pptx_bytes(document_sha256="f" * 64)
    with zipfile.ZipFile(io.BytesIO(body)) as archive:
        slide = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    slide.set("data-document-sha256", DOCUMENT_SHA256)
    return _replace_zip_member(
        body,
        "ppt/slides/slide1.xml",
        ET.tostring(slide, encoding="utf-8", xml_declaration=True),
    )


def _pptx_with_duplicate_slide_identity() -> bytes:
    body = _pptx_bytes()
    with zipfile.ZipFile(io.BytesIO(body)) as archive:
        relationships = ET.fromstring(archive.read("ppt/_rels/presentation.xml.rels"))
    slide_relationship = next(
        node for node in relationships if str(node.get("Type") or "").endswith("/slide")
    )
    relationship_id = slide_relationship.get("Id")
    assert relationship_id
    presentation = f"""<p:presentation xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldIdLst><p:sldId id="256" r:id="{relationship_id}"/><p:sldId id="256" r:id="{relationship_id}"/></p:sldIdLst>
</p:presentation>""".encode()
    return _replace_zip_member(body, "ppt/presentation.xml", presentation)


def _pptx_with_hash_shape_outside_slide() -> bytes:
    presentation = Presentation(io.BytesIO(_pptx_bytes()))
    slide = presentation.slides[0]
    hash_shape = next(
        shape for shape in slide.shapes if hasattr(shape, "text") and DOCUMENT_SHA256 in shape.text
    )
    hash_shape.left = presentation.slide_width
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _pptx_with_zero_sized_hash_shape() -> bytes:
    presentation = Presentation(io.BytesIO(_pptx_bytes()))
    slide = presentation.slides[0]
    hash_shape = next(
        shape for shape in slide.shapes if hasattr(shape, "text") and DOCUMENT_SHA256 in shape.text
    )
    hash_shape.width = 0
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _pptx_with_external_relationship() -> bytes:
    body = _pptx_bytes()
    relationship_path = "ppt/_rels/presentation.xml.rels"
    with zipfile.ZipFile(io.BytesIO(body)) as archive:
        relationships = ET.fromstring(archive.read(relationship_path))
    ET.SubElement(
        relationships,
        "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship",
        {
            "Id": "rIdExternal",
            "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            "Target": "https://external.example.test/classroom",
            "TargetMode": "External",
        },
    )
    return _replace_zip_member(
        body,
        relationship_path,
        ET.tostring(relationships, encoding="utf-8", xml_declaration=True),
    )


def _pptx_with_extra_member(name: str, body: bytes = b"active-content") -> bytes:
    files: dict[str, bytes] = {}
    pptx_body = _pptx_bytes()
    with zipfile.ZipFile(io.BytesIO(pptx_body)) as source:
        for info in source.infolist():
            files[info.filename] = source.read(info)
    files[name] = body
    return _zip_bytes(files)


def _html_bytes(
    *,
    classroom_body: bytes = CLASSROOM_BODY,
    body_document_sha256: str | None = DOCUMENT_SHA256,
    include_main: bool = True,
    close_document: bool = True,
) -> bytes:
    body_attribute = (
        f' data-document-sha256="{body_document_sha256}"'
        if body_document_sha256 is not None
        else ""
    )
    main = "<main><h1>First release classroom</h1></main>" if include_main else ""
    closing = "</body></html>" if close_document else ""
    return (
        '<!doctype html><html><head><meta charset="utf-8"><title>Classroom</title>'
        f"</head><body{body_attribute}>"
        f"{main}"
        '<script id="classroom" type="application/json">'
        f"{classroom_body.decode('utf-8')}"
        f"</script>{closing}"
    ).encode()


def _box(kind: bytes, payload: bytes, *, extended: bool = False) -> bytes:
    if extended:
        return struct.pack(">I4sQ", 1, kind, 16 + len(payload)) + payload
    return struct.pack(">I4s", 8 + len(payload), kind) + payload


def _full_box_header(*, version: int = 0) -> bytes:
    return bytes((version, 0, 0, 0))


def _mvhd_payload(
    *,
    version: int = 0,
    timescale: int = 1_000,
    duration: int = 5_000,
) -> bytes:
    if version == 0:
        header = _full_box_header(version=version) + struct.pack(">IIII", 0, 0, timescale, duration)
        return header + b"\x00" * (100 - len(header))
    if version == 1:
        header = _full_box_header(version=version) + struct.pack(">QQIQ", 0, 0, timescale, duration)
        return header + b"\x00" * (112 - len(header))
    return _full_box_header(version=version) + b"\x00" * 96


def _tkhd_payload(
    *,
    version: int = 0,
    track_id: int = 1,
    duration: int = 5_000,
) -> bytes:
    if version == 0:
        header = _full_box_header(version=version) + struct.pack(
            ">IIIII", 0, 0, track_id, 0, duration
        )
        minimum_bytes = 84
    else:
        header = _full_box_header(version=version) + struct.pack(
            ">QQIIQ", 0, 0, track_id, 0, duration
        )
        minimum_bytes = 96
    geometry = struct.pack(">II", 1_280 << 16, 720 << 16)
    return header + b"\x00" * (minimum_bytes - len(header) - len(geometry)) + geometry


def _mdhd_payload(
    *,
    version: int = 0,
    timescale: int = 1_000,
    duration: int = 5_000,
) -> bytes:
    if version == 0:
        return (
            _full_box_header(version=version)
            + struct.pack(">IIII", 0, 0, timescale, duration)
            + b"\x00\x00\x00\x00"
        )
    return (
        _full_box_header(version=version)
        + struct.pack(">QQIQ", 0, 0, timescale, duration)
        + b"\x00\x00\x00\x00"
    )


def _hdlr_payload(*, handler_type: bytes = b"vide") -> bytes:
    return b"\x00" * 8 + handler_type + b"\x00" * 13


def _visual_sample_entry(*, codec: bytes = b"avc1") -> bytes:
    compressor_name = b"\x00" * 32
    payload = b"".join(
        (
            b"\x00" * 6,
            struct.pack(">H", 1),
            b"\x00" * 16,
            struct.pack(">HH", 1_280, 720),
            struct.pack(">II", 72 << 16, 72 << 16),
            b"\x00" * 4,
            struct.pack(">H", 1),
            compressor_name,
            struct.pack(">HH", 24, 0xFFFF),
        )
    )
    assert len(payload) == 78
    return _box(codec, payload)


def _sample_table_payload(
    *,
    codec: bytes = b"avc1",
    stsd_entry_count: int = 1,
    empty_table: bytes | None = None,
    omit_table: bytes | None = None,
    chunk_offset_box: bytes = b"stco",
) -> bytes:
    stsd_entries = _visual_sample_entry(codec=codec) if stsd_entry_count else b""
    tables = {
        b"stsd": _full_box_header() + struct.pack(">I", stsd_entry_count) + stsd_entries,
        b"stts": _full_box_header() + struct.pack(">III", 1, 1, 1_000),
        b"stsc": _full_box_header() + struct.pack(">IIII", 1, 1, 1, 1),
        b"stsz": _full_box_header() + struct.pack(">III", 0, 1, 128),
        chunk_offset_box: _full_box_header()
        + struct.pack(">I", 1)
        + (struct.pack(">Q", 1) if chunk_offset_box == b"co64" else struct.pack(">I", 1)),
    }
    if empty_table is not None:
        tables[empty_table] = _full_box_header() + (
            struct.pack(">II", 0, 0) if empty_table == b"stsz" else struct.pack(">I", 0)
        )
    if omit_table is not None:
        tables.pop(omit_table)
    return b"".join(_box(kind, payload) for kind, payload in tables.items())


def _mdia_payload(
    *,
    handler_type: bytes = b"vide",
    include_stbl: bool = True,
    mdhd_timescale: int = 1_000,
    mdhd_duration: int = 5_000,
    sample_table_payload: bytes | None = None,
) -> bytes:
    minf_payload = (
        _box(
            b"stbl",
            _sample_table_payload() if sample_table_payload is None else sample_table_payload,
        )
        if include_stbl
        else _box(b"vmhd", b"\x00" * 12)
    )
    return (
        _box(
            b"mdhd",
            _mdhd_payload(timescale=mdhd_timescale, duration=mdhd_duration),
        )
        + _box(b"hdlr", _hdlr_payload(handler_type=handler_type))
        + _box(b"minf", minf_payload)
    )


def _trak_payload(
    *,
    handler_type: bytes = b"vide",
    include_mdia: bool = True,
    include_stbl: bool = True,
    track_id: int = 1,
    tkhd_payload: bytes | None = None,
    mdhd_timescale: int = 1_000,
    mdhd_duration: int = 5_000,
    sample_table_payload: bytes | None = None,
) -> bytes:
    payload = _box(
        b"tkhd",
        _tkhd_payload(track_id=track_id) if tkhd_payload is None else tkhd_payload,
    )
    if include_mdia:
        payload += _box(
            b"mdia",
            _mdia_payload(
                handler_type=handler_type,
                include_stbl=include_stbl,
                mdhd_timescale=mdhd_timescale,
                mdhd_duration=mdhd_duration,
                sample_table_payload=sample_table_payload,
            ),
        )
    return payload


def _moov_payload(
    *,
    handler_type: bytes = b"vide",
    include_mdia: bool = True,
    include_stbl: bool = True,
    mvhd_timescale: int = 1_000,
    mvhd_duration: int = 5_000,
    track_id: int = 1,
    tkhd_payload: bytes | None = None,
    mdhd_timescale: int = 1_000,
    mdhd_duration: int = 5_000,
    sample_table_payload: bytes | None = None,
) -> bytes:
    return _box(
        b"mvhd",
        _mvhd_payload(timescale=mvhd_timescale, duration=mvhd_duration),
    ) + _box(
        b"trak",
        _trak_payload(
            handler_type=handler_type,
            include_mdia=include_mdia,
            include_stbl=include_stbl,
            track_id=track_id,
            tkhd_payload=tkhd_payload,
            mdhd_timescale=mdhd_timescale,
            mdhd_duration=mdhd_duration,
            sample_table_payload=sample_table_payload,
        ),
    )


def _mp4_bytes(*, moov_payload: bytes | None = None) -> bytes:
    return b"".join(
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isommp42"),
            _box(b"moov", _moov_payload() if moov_payload is None else moov_payload),
            _box(b"mdat", b"classroom-video"),
        )
    )


def _write_valid_artifacts(root: Path) -> None:
    root.mkdir(parents=True, exist_ok=True)
    (root / "classroom.zip").write_bytes(_classroom_zip_bytes())
    (root / "classroom.pptx").write_bytes(_pptx_bytes())
    (root / "classroom.html").write_bytes(_html_bytes())
    (root / "classroom.mp4").write_bytes(_mp4_bytes())


def _report(root: Path) -> dict[str, object]:
    exports: dict[str, object] = {}
    for sequence, (kind, relative_path) in enumerate(EXPORT_PATHS.items(), start=1):
        body = (root / relative_path).read_bytes()
        exports[kind] = {
            "exportId": f"export-{sequence}",
            "jobId": f"job-{sequence}",
            "status": "succeeded",
            "progressPercent": 100,
            "downloadReady": True,
            "relativePath": relative_path,
            "contentType": CONTENT_TYPES[kind],
            "contentDisposition": f"attachment; filename*=UTF-8''{relative_path}",
            "byteLength": len(body),
            "sha256": hashlib.sha256(body).hexdigest(),
        }
    return {
        "schemaVersion": 1,
        "producer": "classroom-export-probe",
        "candidate": _candidate(),
        "releaseRun": _release_run(),
        "observedAt": OBSERVED_AT,
        "baseUrl": BASE_URL,
        "tenantId": "tenant-classroom-export",
        "classroomVersionId": CLASSROOM_VERSION_ID,
        "documentSha256": DOCUMENT_SHA256,
        "exports": exports,
    }


def _parse(report: dict[str, object], root: Path) -> dict[str, object]:
    return _module().parse_classroom_export_report(
        _canonical_json(report),
        artifact_root=root,
        candidate=_candidate(),
        release_run=_release_run(),
        expected_base_url=BASE_URL,
    )


def test_classroom_exports_command_record_is_fixed() -> None:
    assert _module().classroom_exports_command_record() == {
        "runner": "python",
        "script": "scripts/classroom_export_probe.py",
        "arguments": ["--profile", "first-release"],
    }


def test_valid_report_opens_all_exports_and_derives_checks(tmp_path: Path) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)

    parsed = _parse(report, tmp_path)

    assert parsed == report
    assert _module().derive_classroom_export_checks(
        _canonical_json(report),
        artifact_root=tmp_path,
        candidate=_candidate(),
        release_run=_release_run(),
        expected_base_url=BASE_URL,
    ) == {
        "zipOpened": True,
        "pptxOpened": True,
        "offlineHtmlOpened": True,
        "mp4Opened": True,
    }


def test_retained_artifact_handles_avoid_path_reopens(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    module = _module()
    handles = {
        kind: (tmp_path / relative_path).open("rb") for kind, relative_path in EXPORT_PATHS.items()
    }

    def fail_path_reopen(*_args: object, **_kwargs: object) -> None:
        pytest.fail("retained artifact validation must not reopen paths")

    monkeypatch.setattr(module, "_open_export_path", fail_path_reopen)
    monkeypatch.setattr(module, "_hash_file", fail_path_reopen)
    try:
        parsed = module.parse_classroom_export_report(
            _canonical_json(report),
            artifact_root=tmp_path,
            artifact_handles=handles,
            candidate=_candidate(),
            release_run=_release_run(),
            expected_base_url=BASE_URL,
        )
        checks = module.derive_classroom_export_checks(
            _canonical_json(report),
            artifact_root=tmp_path,
            artifact_handles=handles,
            candidate=_candidate(),
            release_run=_release_run(),
            expected_base_url=BASE_URL,
        )
    finally:
        for handle in handles.values():
            handle.close()

    assert parsed == report
    assert checks == {
        "zipOpened": True,
        "pptxOpened": True,
        "offlineHtmlOpened": True,
        "mp4Opened": True,
    }


def test_derive_checks_reopens_and_revalidates_raw_artifacts(tmp_path: Path) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    body = _canonical_json(report)
    _parse(report, tmp_path)
    (tmp_path / "classroom.html").write_bytes(b"tampered after parse")

    with pytest.raises(ValueError, match="byte length|SHA-256|HTML"):
        _module().derive_classroom_export_checks(
            body,
            artifact_root=tmp_path,
            candidate=_candidate(),
            release_run=_release_run(),
            expected_base_url=BASE_URL,
        )


def test_report_must_be_canonical_json(tmp_path: Path) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    noncanonical = json.dumps(report, indent=2).encode()

    with pytest.raises(ValueError, match="canonical"):
        _module().parse_classroom_export_report(
            noncanonical,
            artifact_root=tmp_path,
            candidate=_candidate(),
            release_run=_release_run(),
            expected_base_url=BASE_URL,
        )


@pytest.mark.parametrize("binding", ("candidate", "releaseRun", "baseUrl"))
def test_report_is_bound_to_supplied_release_material(tmp_path: Path, binding: str) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    if binding == "candidate":
        report[binding] = {**_candidate(), "sourceHead": "b" * 40}
    elif binding == "releaseRun":
        report[binding] = {**_release_run(), "runId": "different-run"}
    else:
        report[binding] = "https://other.example.test"

    with pytest.raises(ValueError, match="binding|URL"):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    "observed_at",
    (
        None,
        True,
        "",
        "2026-08-28T12:34:56",
        "2026-08-28T12:34:56+00:00",
        "2026-08-28T12:34:56.123z",
        "2026-02-30T12:34:56Z",
    ),
)
def test_report_requires_a_real_utc_z_observation_time(tmp_path: Path, observed_at: object) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    report["observedAt"] = observed_at

    with pytest.raises(ValueError, match="timestamp"):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("mutation", "message"),
    (
        (lambda report: report.update(extra=True), "invalid"),
        (lambda report: report.update(schemaVersion=True), "schema"),
        (lambda report: report.update(producer="other-probe"), "producer"),
        (lambda report: report.update(tenantId=""), "tenant"),
        (lambda report: report.update(classroomVersionId="mutable version"), "version"),
        (lambda report: report.update(documentSha256="0" * 64), "document"),
        (lambda report: report.pop("observedAt"), "invalid"),
        (lambda report: report["exports"].pop("mp4"), "exports"),
        (lambda report: report["exports"].update(extra={}), "exports"),
    ),
)
def test_report_requires_the_exact_schema(tmp_path: Path, mutation, message: str) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    mutation(report)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("field", "value", "message"),
    (
        ("exportId", "", "identity"),
        ("jobId", "job with space", "identity"),
        ("status", "running", "status"),
        ("progressPercent", True, "progress"),
        ("progressPercent", 99, "progress"),
        ("downloadReady", 1, "download"),
        ("contentType", "application/octet-stream", "content type"),
        ("contentDisposition", "inline", "disposition"),
        (
            "contentDisposition",
            'attachment; filename="classroom.zip"',
            "disposition",
        ),
        ("byteLength", True, "byte length"),
        ("sha256", "a" * 63, "SHA-256"),
    ),
)
def test_export_metadata_is_strict(tmp_path: Path, field: str, value: object, message: str) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    report["exports"]["classroom_zip"][field] = value

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    "relative_path",
    (
        "exports/classroom.zip",
        "../classroom.zip",
        "/tmp/classroom.zip",
        "C:\\outside\\classroom.zip",
    ),
)
def test_export_paths_are_fixed_safe_basenames(tmp_path: Path, relative_path: str) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    report["exports"]["classroom_zip"]["relativePath"] = relative_path

    with pytest.raises(ValueError, match="path"):
        _parse(report, tmp_path)


def test_export_path_must_not_escape_through_a_symlink(tmp_path: Path) -> None:
    artifact_root = tmp_path / "raw"
    _write_valid_artifacts(artifact_root)
    report = _report(artifact_root)
    original = artifact_root / "classroom.zip"
    retained = artifact_root / "classroom.zip.retained"
    original.rename(retained)
    external = tmp_path / "outside.zip"
    external.write_bytes(retained.read_bytes())
    try:
        original.symlink_to(external)
    except OSError as exc:
        pytest.skip(f"file symlinks are unavailable: {exc}")

    with pytest.raises(ValueError, match="path|root|symlink"):
        _parse(report, artifact_root)


def test_export_basename_must_not_be_a_symlink_within_artifact_root(
    tmp_path: Path,
) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    original = tmp_path / "classroom.zip"
    retained = tmp_path / "classroom.zip.retained"
    original.rename(retained)
    try:
        original.symlink_to(retained)
    except OSError as exc:
        pytest.skip(f"file symlinks are unavailable: {exc}")

    with pytest.raises(ValueError, match="symlink"):
        _parse(report, tmp_path)


def test_artifact_root_must_not_be_a_symlink(tmp_path: Path) -> None:
    actual_root = tmp_path / "actual"
    _write_valid_artifacts(actual_root)
    report = _report(actual_root)
    linked_root = tmp_path / "linked"
    try:
        linked_root.symlink_to(actual_root, target_is_directory=True)
    except OSError as exc:
        pytest.skip(f"directory symlinks are unavailable: {exc}")

    with pytest.raises(ValueError, match="symlink|junction"):
        _parse(report, linked_root)


@pytest.mark.parametrize("field", ("byteLength", "sha256"))
def test_export_bytes_and_hash_are_recomputed(tmp_path: Path, field: str) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    export = report["exports"]["classroom_zip"]
    export[field] = export[field] + 1 if field == "byteLength" else "f" * 64

    with pytest.raises(ValueError, match="byte length|SHA-256"):
        _parse(report, tmp_path)


def test_total_artifact_size_is_bounded(tmp_path: Path, monkeypatch) -> None:
    _write_valid_artifacts(tmp_path)
    report = _report(tmp_path)
    total = sum(item["byteLength"] for item in report["exports"].values())
    monkeypatch.setattr(_module(), "MAX_TOTAL_EXPORT_BYTES", total - 1)

    with pytest.raises(ValueError, match="total.*large"):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("zip_body", "message"),
    (
        (
            _zip_bytes({"classroom.json": CLASSROOM_BODY}),
            "members|media",
        ),
        (
            _classroom_zip_bytes(classroom_body=CLASSROOM_BODY + b" "),
            "classroom.json.*hash",
        ),
        (
            _classroom_zip_bytes(media_files={"media/voice.mp3": b"wrong-media"}),
            "media.*size|media.*hash",
        ),
        (
            _classroom_zip_bytes(extra_files={"extra.txt": b"extra"}),
            "members",
        ),
        (
            _classroom_zip_bytes(extra_files={"../escape.txt": b"escape"}),
            "unsafe",
        ),
    ),
)
def test_classroom_zip_must_be_safe_and_bound(
    tmp_path: Path, zip_body: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.zip").write_bytes(zip_body)
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


@pytest.mark.parametrize("kind", ("classroom_zip", "pptx"))
def test_archive_forbidden_byte_scan_detects_cross_chunk_secret(kind: str) -> None:
    module = _module()
    token = b"cross-chunk-live-fixture-token"
    member_body = b"x" * (1024 * 1024 - len(token) // 2) + token + b"tail"
    archive_body = io.BytesIO()
    with zipfile.ZipFile(archive_body, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("payload.bin", member_body)
    assert token not in archive_body.getvalue()

    archive_body.seek(0)
    assert module.classroom_export_archive_contains_forbidden_bytes(
        archive_body,
        kind=kind,
        forbidden={token},
    )
    assert archive_body.tell() == 0


def test_classroom_zip_rejects_invalid_classroom_document_even_when_hash_matches(
    tmp_path: Path,
) -> None:
    _write_valid_artifacts(tmp_path)
    invalid_document = b'{"not":"a-classroom-document"}'
    (tmp_path / "classroom.zip").write_bytes(_classroom_zip_bytes(classroom_body=invalid_document))
    report = _report(tmp_path)
    report["documentSha256"] = hashlib.sha256(invalid_document).hexdigest()

    with pytest.raises(ValueError, match="ClassroomDocument|classroom.json"):
        _parse(report, tmp_path)


def test_classroom_zip_rejects_a_different_classroom_version_with_coherent_hash(
    tmp_path: Path,
) -> None:
    _write_valid_artifacts(tmp_path)
    other_document = _valid_classroom_body(classroom_version_id="other-version")
    (tmp_path / "classroom.zip").write_bytes(_classroom_zip_bytes(classroom_body=other_document))
    report = _report(tmp_path)
    report["documentSha256"] = hashlib.sha256(other_document).hexdigest()

    with pytest.raises(ValueError, match="version"):
        _parse(report, tmp_path)


def test_classroom_zip_rejects_corrupt_archive(tmp_path: Path) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.zip").write_bytes(b"PK\x03\x04truncated")
    report = _report(tmp_path)

    with pytest.raises(ValueError, match="ZIP|archive"):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("pptx_body", "message"),
    (
        (b"PK\x03\x04truncated", "PPTX|archive"),
        (
            _zip_bytes(
                {
                    "[Content_Types].xml": b"<Types xmlns='http://schemas.openxmlformats.org/package/2006/content-types'/>",
                    "_rels/.rels": b"<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'/>",
                }
            ),
            "presentation|slide",
        ),
        (_pptx_bytes(document_sha256="f" * 64), "document"),
        (_pptx_with_hash_only_in_xml_attribute(), "visible|document"),
        (
            _replace_zip_member(
                _pptx_bytes(),
                "[Content_Types].xml",
                b"""<Wrong xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
</Wrong>""",
            ),
            "content types",
        ),
        (
            _pptx_with_duplicate_slide_identity(),
            "slide identity",
        ),
    ),
    ids=(
        "truncated-archive",
        "missing-presentation",
        "wrong-visible-hash",
        "hash-only-in-xml-attribute",
        "invalid-content-types-root",
        "duplicate-slide-identity",
    ),
)
def test_pptx_must_be_valid_ooxml_with_a_bound_slide(
    tmp_path: Path, pptx_body: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.pptx").write_bytes(pptx_body)
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("pptx_body", "message"),
    (
        (_pptx_with_zero_sized_hash_shape(), "visible|bounds|shape"),
        (_pptx_with_hash_shape_outside_slide(), "visible|bounds|shape"),
        (_pptx_with_external_relationship(), "external relationship"),
        (_pptx_with_extra_member("ppt/vbaProject.bin"), "macro|active content"),
        (
            _pptx_with_extra_member("ppt/activeX/activeX1.bin"),
            "ActiveX|active content",
        ),
        (
            _pptx_with_extra_member("ppt/embeddings/oleObject1.bin"),
            "OLE|active content",
        ),
    ),
    ids=(
        "zero-sized-hash-shape",
        "out-of-bounds-hash-shape",
        "external-relationship",
        "vba-macro",
        "activex",
        "ole-object",
    ),
)
def test_pptx_rejects_nonvisible_hash_and_active_package_content(
    tmp_path: Path, pptx_body: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.pptx").write_bytes(pptx_body)
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("html_body", "message"),
    (
        (
            b"<!doctype html><html><head><title>Classroom</title></head>"
            b"<body><main>Missing script</main></body></html>",
            "script",
        ),
        (_html_bytes(classroom_body=b'{"invalid":true}'), "ClassroomDocument|document"),
        (_html_bytes(include_main=False), "main"),
        (_html_bytes(close_document=False), "closed|structure|HTML"),
        (_html_bytes(body_document_sha256="f" * 64), "document"),
        (
            b"<!doctype html><html><body><main></main></body></html>trailing",
            "structure|HTML",
        ),
        (
            b"<!doctype html><html><body><main></body></main></html>",
            "structure|HTML",
        ),
        (b"\xff\xfe", "UTF-8|HTML"),
    ),
)
def test_offline_html_must_bind_the_document_and_contain_main(
    tmp_path: Path, html_body: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.html").write_bytes(html_body)
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


def test_offline_html_accepts_production_without_optional_body_hash(
    tmp_path: Path,
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.html").write_bytes(_html_bytes(body_document_sha256=None))
    report = _report(tmp_path)

    assert _parse(report, tmp_path) == report


def test_offline_html_rejects_a_coherent_different_classroom_version(
    tmp_path: Path,
) -> None:
    other_document = _valid_classroom_body(classroom_version_id="other-version")
    other_sha256 = hashlib.sha256(other_document).hexdigest()
    html_path = tmp_path / "classroom.html"
    html_path.write_bytes(
        _html_bytes(
            classroom_body=other_document,
            body_document_sha256=other_sha256,
        )
    )

    with pytest.raises(ValueError, match="version"):
        _module()._validate_offline_html(
            html_path,
            classroom_version_id=CLASSROOM_VERSION_ID,
            document_sha256=other_sha256,
        )


@pytest.mark.parametrize(
    ("html_body", "message"),
    (
        (_html_bytes().replace(b"<!doctype html>", b""), "doctype"),
        (
            _html_bytes().replace(b"<!doctype html>", b"<!doctype svg>"),
            "doctype",
        ),
        (
            _html_bytes().replace(
                b'<head><meta charset="utf-8"><title>Classroom</title></head>',
                b"",
            ),
            "head|structure",
        ),
        (
            _html_bytes()
            .replace(b"<main>", b"<section><main>")
            .replace(b"</main>", b"</main></section>"),
            "main|structure",
        ),
        (
            _html_bytes().replace(b"</main>", b"</main><main></main>"),
            "main|structure",
        ),
    ),
    ids=(
        "missing-doctype",
        "wrong-doctype",
        "missing-head",
        "nested-main",
        "duplicate-main",
    ),
)
def test_offline_html_requires_strict_html5_document_structure(
    tmp_path: Path, html_body: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.html").write_bytes(html_body)
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("html_body", "message"),
    (
        (
            _html_bytes().replace(
                b"<main>",
                b'<main><img src="data:image/png;base64,AA==">',
            ),
            "src|resource",
        ),
        (
            _html_bytes().replace(
                b"<main>",
                b'<main><source srcset="image.png 1x">',
            ),
            "srcset|resource",
        ),
        (
            _html_bytes().replace(b"<main>", b'<main><a href="#local">link</a>'),
            "href|resource",
        ),
        (
            _html_bytes().replace(b"<main>", b'<main><form action="/submit"></form>'),
            "form|action|active",
        ),
        (
            _html_bytes().replace(b"<main>", b"<main><iframe></iframe>"),
            "iframe|active",
        ),
        (
            _html_bytes().replace(
                b"</main>",
                b"</main><script>window.location='//external.example.test'</script>",
            ),
            "script|active",
        ),
        (
            _html_bytes().replace(
                b"<main>",
                b"<main><button onclick=\"window.location='//external.example.test'\">",
            ),
            "event|onclick|active",
        ),
        (
            _html_bytes().replace(
                b"</title>",
                b'</title><meta http-equiv="refresh" content="0;url=//external.example.test">',
            ),
            "refresh|active",
        ),
        (
            _html_bytes().replace(
                b"<main>",
                b'<main><div data-location="//external.example.test">',
            ),
            "external|protocol-relative",
        ),
        (
            _html_bytes().replace(
                b"</title>",
                b'</title><style>@import "theme.css";</style>',
            ),
            "CSS|import|resource",
        ),
        (
            _html_bytes().replace(
                b"</title>",
                b"</title><style>main{background:url(data:image/png;base64,AA==)}</style>",
            ),
            "CSS|url|resource",
        ),
        (
            _html_bytes().replace(
                b"<main>",
                b'<main><div style="background:url(//external.example.test/a.png)">',
            ),
            "CSS|url|resource",
        ),
    ),
    ids=(
        "src-attribute",
        "srcset-attribute",
        "href-attribute",
        "form-action",
        "iframe",
        "non-classroom-script",
        "event-handler-attribute",
        "meta-refresh",
        "protocol-relative-url",
        "css-import",
        "css-url-rule",
        "css-url-style-attribute",
    ),
)
def test_offline_html_rejects_active_and_external_resources(
    tmp_path: Path, html_body: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.html").write_bytes(html_body)
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


@pytest.mark.parametrize(
    ("mp4_body", "message"),
    (
        (_box(b"ftyp", b"isom\0\0\0\0") + _box(b"mdat", b"data"), "moov"),
        (struct.pack(">I4s", 0, b"ftyp") + b"isom", "size"),
        (struct.pack(">I4s", 32, b"ftyp") + b"short", "truncated"),
        (struct.pack(">I4sQ", 1, b"ftyp", 12) + b"data", "extended|size"),
        (struct.pack(">I4s", 4, b"ftyp"), "size"),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00x") + _box(b"moov", b"") + _box(b"mdat", b"data"),
            "ftyp",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(b"moov", b"")
            + _box(b"mdat", b"data"),
            "moov",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(b"moov", b"truncated")
            + _box(b"mdat", b"data"),
            "moov|truncated|structure",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(
                b"moov",
                _box(b"mvhd", _mvhd_payload(version=2)) + _box(b"trak", _trak_payload()),
            )
            + _box(b"mdat", b"data"),
            "mvhd.*version",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(
                b"moov",
                _box(b"mvhd", b"\x00" * 16) + _box(b"trak", _trak_payload()),
            )
            + _box(b"mdat", b"data"),
            "mvhd.*short",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(b"moov", _box(b"mvhd", _mvhd_payload()))
            + _box(b"mdat", b"data"),
            "trak|moov",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(b"moov", _moov_payload(include_mdia=False))
            + _box(b"mdat", b"data"),
            "mdia",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(b"moov", _moov_payload(handler_type=b"soun"))
            + _box(b"mdat", b"data"),
            "video",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(b"moov", _moov_payload(include_stbl=False))
            + _box(b"mdat", b"data"),
            "stbl",
        ),
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isom")
            + _box(b"moov", _moov_payload())
            + _box(b"mdat", b""),
            "mdat",
        ),
    ),
)
def test_mp4_requires_complete_top_level_boxes(
    tmp_path: Path, mp4_body: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.mp4").write_bytes(mp4_body)
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


def test_mp4_accepts_a_well_formed_extended_size_box(tmp_path: Path) -> None:
    _write_valid_artifacts(tmp_path)
    body = b"".join(
        (
            _box(b"ftyp", b"isom\x00\x00\x02\x00isommp42", extended=True),
            _box(b"moov", _moov_payload()),
            _box(b"mdat", b"classroom-video"),
        )
    )
    (tmp_path / "classroom.mp4").write_bytes(body)
    report = _report(tmp_path)

    assert _parse(report, tmp_path) == report


@pytest.mark.parametrize(
    ("moov_payload", "message"),
    (
        (_moov_payload(mvhd_timescale=0), "mvhd.*timescale"),
        (_moov_payload(mvhd_duration=0), "mvhd.*duration"),
        (_moov_payload(mdhd_timescale=0), "mdhd.*timescale"),
        (_moov_payload(mdhd_duration=0), "mdhd.*duration"),
        (_moov_payload(track_id=0), "tkhd.*track"),
        (_moov_payload(tkhd_payload=b"\x00" * 16), "tkhd.*short|tkhd.*length"),
        (
            _moov_payload(sample_table_payload=_sample_table_payload(stsd_entry_count=0)),
            "stsd.*entry|sample entry",
        ),
        (
            _moov_payload(sample_table_payload=_sample_table_payload(codec=b"mp4a")),
            "stsd.*visual|sample entry",
        ),
        (
            _moov_payload(sample_table_payload=_sample_table_payload(empty_table=b"stts")),
            "stts.*empty|stts.*entry",
        ),
        (
            _moov_payload(sample_table_payload=_sample_table_payload(empty_table=b"stsc")),
            "stsc.*empty|stsc.*entry",
        ),
        (
            _moov_payload(sample_table_payload=_sample_table_payload(empty_table=b"stsz")),
            "stsz.*empty|stsz.*sample",
        ),
        (
            _moov_payload(sample_table_payload=_sample_table_payload(empty_table=b"stco")),
            "stco.*empty|chunk offset",
        ),
        (
            _moov_payload(sample_table_payload=_sample_table_payload(omit_table=b"stco")),
            "stco|co64|chunk offset",
        ),
    ),
    ids=(
        "zero-movie-timescale",
        "zero-movie-duration",
        "zero-media-timescale",
        "zero-media-duration",
        "zero-track-id",
        "short-track-header",
        "empty-sample-description",
        "unsupported-visual-codec",
        "empty-time-to-sample",
        "empty-sample-to-chunk",
        "empty-sample-size",
        "empty-32-bit-chunk-offset",
        "missing-chunk-offset-table",
    ),
)
def test_mp4_requires_nonzero_timing_and_complete_video_sample_tables(
    tmp_path: Path, moov_payload: bytes, message: str
) -> None:
    _write_valid_artifacts(tmp_path)
    (tmp_path / "classroom.mp4").write_bytes(_mp4_bytes(moov_payload=moov_payload))
    report = _report(tmp_path)

    with pytest.raises(ValueError, match=message):
        _parse(report, tmp_path)


def test_mp4_accepts_a_complete_64_bit_chunk_offset_table(tmp_path: Path) -> None:
    _write_valid_artifacts(tmp_path)
    moov_payload = _moov_payload(
        sample_table_payload=_sample_table_payload(chunk_offset_box=b"co64")
    )
    (tmp_path / "classroom.mp4").write_bytes(_mp4_bytes(moov_payload=moov_payload))
    report = _report(tmp_path)

    assert _parse(report, tmp_path) == report
