"""Deterministic acceptance flow for teacher classroom publishing."""

from __future__ import annotations

from dataclasses import dataclass, replace
from datetime import datetime, timezone
import hashlib
from io import BytesIO
import json
from types import SimpleNamespace
from xml.etree import ElementTree
import zipfile

from pptx import Presentation
import pytest

from deeptutor.api.routers.classroom_exports import download_classroom_export
from deeptutor.multi_user.knowledge_access import AuthorizedKnowledgeSource
from deeptutor.services.rag.retrieval_view import stamp_retrieval_view_signature
from deeptutor.teaching.brief_builder import TeachingBriefBuilder
from deeptutor.teaching.contracts import (
    ClassroomDocument,
    OutlineBundle,
    canonical_json_bytes,
    canonical_outline_sha256,
)
from deeptutor.teaching.models.classrooms import (
    InvalidClassroomTransition,
    transition,
)
from deeptutor.teaching.permissions import permissions_for_roles
from deeptutor.teaching.repositories.sources import (
    BoundSourceRecord,
    SavedSourceSnapshot,
)
from deeptutor.teaching.services.classrooms import ClassroomService
from deeptutor.teaching.services.exports import (
    ClassroomExportService,
    ExportRecord,
    ExportSource,
)
from deeptutor.teaching.services.publications import (
    AssignmentRecord,
    ConfirmedPublicationMaterialization,
    MaterializedPublicationArtifact,
    PublicationAccessDenied,
    PublicationMaterializationPlan,
    PublicationService,
    PublicationTarget,
    PublishedVersionRecord,
    TenantPublicationCandidate,
    TenantPublicationItem,
    VersionTarget,
    publication_media_manifest_sha256,
)
from deeptutor.teaching.services.reviews import (
    ReviewAccessDenied,
    ReviewConflict,
    ReviewPolicy,
    ReviewRecord,
    ReviewService,
    ReviewTarget,
)
from deeptutor.teaching.source_snapshots import (
    SourceSnapshotBuilder,
)
from deeptutor.teaching.tenant_context import TenantContext
from tests.teaching.test_classroom_export_service import (
    _Materializer as ExportMemoryMaterializer,
)
from tests.teaching.test_classroom_service import (
    _Generation as ClassroomMemoryGeneration,
)
from tests.teaching.test_classroom_service import (
    _Repository as ClassroomMemoryRepository,
)
from tests.teaching.test_classroom_service import _request as classroom_request
from tests.teaching.test_contracts import valid_classroom_document

NOW = datetime(2026, 8, 9, 8, tzinfo=timezone.utc)
KB_GENERATION_ID = "aaaaaaaa-aaaa-4aaa-8aaa-aaaaaaaaaaaa"
KB_SOURCE_ID = f"user:kb:{KB_GENERATION_ID}"


def _context(user_id: str, roles: set[str]) -> TenantContext:
    return TenantContext(
        tenant_id="tenant-a",
        schema_name="tenant_tenant_a",
        user_id=user_id,
        permissions=permissions_for_roles(
            roles,
            scope_type="class",
            scope_id="class-a",
            tenant_id="tenant-a",
        ),
    )


class _BytesStore:
    tenant_id = "tenant-a"

    def __init__(self) -> None:
        self.content: dict[str, bytes] = {}

    async def open(self, key: str):
        payload = self.content[key]

        async def body():
            yield payload

        return body()


class _StoreProvider:
    def __init__(self, store: _BytesStore) -> None:
        self.store = store

    async def store_for_tenant(self, tenant_id: str) -> _BytesStore:
        assert tenant_id == self.store.tenant_id
        return self.store


class _SourceRepository:
    def __init__(self, pdf_key: str, pdf_sha256: str) -> None:
        self.kb = BoundSourceRecord(
            binding_id="binding-kb",
            snapshot_id="bound-kb",
            source_type="knowledge_base",
            source_id=KB_SOURCE_ID,
            resource_owner_id="teacher-a",
            source_revision="kb-revision-1",
            content_sha256="1" * 64,
            permission_sha256="2" * 64,
            display_name=None,
            upload_id=None,
            object_key=None,
        )
        self.pdf = BoundSourceRecord(
            binding_id="binding-pdf",
            snapshot_id="bound-pdf",
            source_type="pdf",
            source_id="pdf-source-1",
            resource_owner_id="teacher-a",
            source_revision="pdf-revision-1",
            content_sha256=pdf_sha256,
            permission_sha256="3" * 64,
            display_name="motion.pdf",
            upload_id="upload-pdf-1",
            object_key=pdf_key,
        )
        self.persisted: list[tuple[BoundSourceRecord, object, dict[str, object]]] = []

    async def require_authorized_source(self, **kwargs) -> BoundSourceRecord:
        return self.pdf if kwargs["source_type"] == "pdf" else self.kb

    async def persist_authorized_snapshot(self, bound, snapshot, **kwargs):
        self.persisted.append((bound, snapshot, kwargs))
        return SavedSourceSnapshot(snapshot_id=snapshot.snapshot_id, created_at=NOW)


class _RagService:
    async def search(self, query: str, kb_name: str) -> dict[str, object]:
        assert "motion" in query.lower()
        assert kb_name == "mechanics"
        result: dict[str, object] = {
            "provider": "llamaindex",
            "index_signature": "mechanics-index-v1",
            "sources": [
                {
                    "chunk_id": "motion-1",
                    "content": "A net force changes an object's momentum.",
                    "source": "mechanics.pdf",
                    "page": 4,
                    "title": "Newton's second law",
                }
            ],
        }
        stamp_retrieval_view_signature(result)
        return result


def _knowledge_source(_reference: str) -> AuthorizedKnowledgeSource:
    return AuthorizedKnowledgeSource(
        resource_id=KB_SOURCE_ID,
        generation_id=KB_GENERATION_ID,
        name="mechanics",
        source="user",
        resource_owner_id="teacher-a",
        read_only=True,
        retrieval_provider="llamaindex",
    )


async def _extract_pdf(_path: str, filename: str) -> str:
    assert filename == "motion.pdf"
    return "--- Page 1 ---\nA force accelerates an object and changes its motion."


def _generated_document(
    asset_id: str,
    version_id: str,
    source_refs: list[dict[str, object]],
) -> dict[str, object]:
    payload = valid_classroom_document()
    payload["classroom_id"] = asset_id
    payload["classroom_version_id"] = version_id
    payload["content_mode"] = "source_grounded"
    payload["open_creation"] = False
    payload["source_refs"] = source_refs
    payload["media_manifest"] = []
    payload["export_manifest"] = []
    openmaic = payload["openmaic"]
    assert isinstance(openmaic, dict)
    scenes = openmaic["scenes"]
    assert isinstance(scenes, list) and isinstance(scenes[0], dict)
    scenes[0]["title"] = "Generated motion lesson"
    scenes[0]["content"] = {
        "type": "slide",
        "canvas": {"text": "Force changes motion."},
    }
    provisional = ClassroomDocument.model_validate(payload)
    unhashed = provisional.model_dump(mode="json", by_alias=True, exclude_none=True)
    unhashed.pop("fileSha256")
    payload["file_sha256"] = hashlib.sha256(
        canonical_json_bytes(unhashed)
    ).hexdigest()
    return ClassroomDocument.model_validate(payload).model_dump(
        mode="json",
        by_alias=True,
        exclude_none=True,
    )


class _FlowClassroomRepository(ClassroomMemoryRepository):
    def __init__(self) -> None:
        super().__init__()
        self.reconciled_assets: list[str] = []

    async def mark_generation_succeeded(self, asset_id: str, job_id: str):
        record = self.records[asset_id]
        assert record.job_id == job_id
        if record.lifecycle_state == "generating_content":
            record = replace(
                record,
                lifecycle_state=transition("generating_content", "editing"),
                status="succeeded",
            )
            self.records[asset_id] = record
            self.reconciled_assets.append(asset_id)
        return record


class _FlowGeneration(ClassroomMemoryGeneration):
    def __init__(self) -> None:
        super().__init__()
        self.stages: dict[str, object] = {}

    async def start_outline(self, **kwargs):
        stage = await super().start_outline(**kwargs)
        job_id = "job-" + hashlib.sha256(kwargs["asset_id"].encode()).hexdigest()[:32]
        stage = replace(stage, job_id=job_id)
        self.stages[job_id] = stage
        return stage

    async def start_content(self, **kwargs):
        stage = await super().start_content(**kwargs)
        stage = replace(stage, job_id=kwargs["job_id"])
        self.stages[stage.job_id] = stage
        return stage

    async def get_stage(self, *, context, job_id):
        del context
        return self.stages[job_id]

    async def complete(self, repository: _FlowClassroomRepository, asset_id: str) -> None:
        record = repository.records[asset_id]
        assert record.job_id is not None
        assert record.teaching_brief is not None
        source_refs = [
            item.model_dump(mode="json", by_alias=True, exclude_none=True)
            for item in record.teaching_brief.source_refs
        ]
        version_id = "generated-" + hashlib.sha256(asset_id.encode()).hexdigest()[:24]
        repository.records[asset_id] = replace(
            record,
            classroom_version_id=version_id,
            document=_generated_document(asset_id, version_id, source_refs),
        )
        self.stages[record.job_id] = replace(
            self.stages[record.job_id],
            status="succeeded",
            classroom_version_id=version_id,
        )


@dataclass(frozen=True, slots=True)
class _LedgerVersion:
    record: PublishedVersionRecord
    tenant_id: str
    owner_id: str
    course_id: str
    class_id: str
    title: str
    review_id: str
    draft_revision: int
    published_by: str
    document: bytes
    media_manifest_sha256: str


class _FlowLedger:
    def __init__(
        self,
        classrooms: _FlowClassroomRepository,
        *,
        policy: ReviewPolicy | None = None,
    ) -> None:
        self.classrooms = classrooms
        self.policy = policy or ReviewPolicy()
        self.reviews: dict[str, ReviewRecord] = {}
        self.review_keys: dict[str, str] = {}
        self.versions: dict[str, _LedgerVersion] = {}
        self.publications_by_review: dict[str, str] = {}
        self.assignments: dict[str, AssignmentRecord] = {}

    async def classroom(self, asset_id: str):
        return await self.classrooms.get_workflow(asset_id)


class _LedgerReviewRepository:
    def __init__(self, ledger: _FlowLedger) -> None:
        self.ledger = ledger

    async def get_policy(self) -> ReviewPolicy:
        return self.ledger.policy

    async def get_target(self, asset_id: str) -> ReviewTarget | None:
        record = await self.ledger.classroom(asset_id)
        if record is None:
            return None
        return ReviewTarget(
            tenant_id=record.tenant_id,
            asset_id=record.asset_id,
            owner_id=record.owner_id,
            course_id=record.course_id,
            class_id=record.class_id,
        )

    async def submit(self, command) -> ReviewRecord:
        existing_id = self.ledger.review_keys.get(command.idempotency_key)
        if existing_id is not None:
            return self.ledger.reviews[existing_id]
        current = await self.ledger.classroom(command.asset_id)
        if (
            current is None
            or current.tenant_id != command.tenant_id
            or current.lifecycle_state != "editing"
            or current.validation_report is None
        ):
            raise ReviewConflict("validated classroom draft is unavailable")
        document = canonical_json_bytes(ClassroomDocument.model_validate(current.document))
        document_sha256 = hashlib.sha256(document).hexdigest()
        report = current.validation_report
        if (
            report.get("draftRevision") != current.revision
            or report.get("documentSha256") != document_sha256
        ):
            raise ReviewConflict("classroom validation binding is stale")
        review_id = (
            "review-"
            + hashlib.sha256(
                f"{command.tenant_id}\0{command.idempotency_key}".encode()
            ).hexdigest()[:24]
        )
        review = ReviewRecord(
            id=review_id,
            tenant_id=current.tenant_id,
            asset_id=current.asset_id,
            draft_id=current.draft_id,
            draft_revision=current.revision,
            document_sha256=document_sha256,
            validation_report_sha256=hashlib.sha256(
                canonical_json_bytes(report)
            ).hexdigest(),
            submitted_by=command.actor_id,
            scope=command.scope,
            class_id=command.class_id,
            status="pending",
            warnings=(),
            reviewer_id=None,
            comment=None,
        )
        self.ledger.reviews[review.id] = review
        self.ledger.review_keys[command.idempotency_key] = review.id
        self.ledger.classrooms.records[current.asset_id] = replace(
            current,
            lifecycle_state=transition("editing", "submitted"),
        )
        return review

    async def list_pending(self) -> tuple[ReviewRecord, ...]:
        return tuple(review for review in self.ledger.reviews.values() if review.status == "pending")

    async def get_review(self, review_id: str) -> ReviewRecord | None:
        return self.ledger.reviews.get(review_id)

    async def get_detail(self, review_id: str):
        del review_id
        return None

    async def decide(self, command) -> ReviewRecord:
        review = self.ledger.reviews.get(command.review_id)
        if review is None or review.status != "pending":
            raise ReviewConflict("review was already decided")
        current = await self.ledger.classroom(review.asset_id)
        if current is None or current.lifecycle_state != "submitted":
            raise ReviewConflict("submitted classroom binding is stale")
        decided = replace(
            review,
            status=command.decision,
            reviewer_id=command.actor_id,
            comment=command.comment,
        )
        self.ledger.reviews[review.id] = decided
        lifecycle_state = transition(
            "submitted",
            "approved" if command.decision == "approved" else "rejected",
        )
        self.ledger.classrooms.records[current.asset_id] = replace(
            current,
            lifecycle_state=lifecycle_state,
        )
        return decided


class _LedgerPublicationMaterializer:
    def __init__(self) -> None:
        self.plans: list[PublicationMaterializationPlan] = []

    async def materialize(
        self,
        plan: PublicationMaterializationPlan,
    ) -> ConfirmedPublicationMaterialization:
        # PostgreSQL triggers/repositories and the production materializer stay in
        # focused integration coverage, unavailable in this no-Docker E2E gate.
        assert hashlib.sha256(plan.document).hexdigest() == plan.document_sha256
        assert publication_media_manifest_sha256(plan.document) == plan.media_manifest_sha256
        self.plans.append(plan)
        return ConfirmedPublicationMaterialization(
            manifest_sha256=plan.manifest_sha256,
            media_manifest_sha256=plan.media_manifest_sha256,
            artifacts=(
                MaterializedPublicationArtifact(
                    relative_name="classroom.json",
                    object_key=(
                        f"tenants/{plan.tenant_id}/classrooms/{plan.asset_id}/"
                        f"versions/{plan.version_number}/classroom.json"
                    ),
                    sha256=plan.document_sha256,
                    size_bytes=len(plan.document),
                    mime_type="application/json",
                    artifact_kind="dsl_json",
                    media_id=None,
                ),
            ),
        )


class _LedgerPublicationRepository:
    def __init__(self, ledger: _FlowLedger) -> None:
        self.ledger = ledger

    async def get_policy(self) -> ReviewPolicy:
        return self.ledger.policy

    def _review_for_asset(self, asset_id: str) -> ReviewRecord | None:
        return next(
            (
                review
                for review in reversed(tuple(self.ledger.reviews.values()))
                if review.asset_id == asset_id
            ),
            None,
        )

    async def get_publication_target(self, asset_id: str) -> PublicationTarget | None:
        review = self._review_for_asset(asset_id)
        record = await self.ledger.classroom(asset_id)
        if review is None or record is None:
            return None
        return PublicationTarget(
            tenant_id=record.tenant_id,
            asset_id=record.asset_id,
            owner_id=record.owner_id,
            course_id=record.course_id,
            class_id=record.class_id,
            review_id=review.id,
            review_scope=review.scope,
            review_status=review.status,
            submitted_by=review.submitted_by,
            draft_revision=review.draft_revision,
            document_sha256=review.document_sha256,
        )

    async def list_tenant_library(self):
        items = tuple(
            TenantPublicationItem(
                tenant_id=version.tenant_id,
                publication_id=f"publication-{version.record.version_id}",
                version_id=version.record.version_id,
                asset_id=version.record.asset_id,
                version_number=version.record.version_number,
                title=version.title,
                course_id=version.course_id,
                document_sha256=version.record.document_sha256,
                published_by=version.published_by,
                created_at=NOW,
                scope=version.record.publication_scope,
            )
            for version in self.ledger.versions.values()
            if version.record.publication_scope == "tenant"
        )
        candidates: list[TenantPublicationCandidate] = []
        for review in self.ledger.reviews.values():
            if (
                review.status != "approved"
                or review.scope != "tenant"
                or review.id in self.ledger.publications_by_review
            ):
                continue
            record = await self.ledger.classroom(review.asset_id)
            if record is None:
                continue
            candidates.append(
                TenantPublicationCandidate(
                    tenant_id=review.tenant_id,
                    review_id=review.id,
                    asset_id=review.asset_id,
                    title=record.title,
                    course_id=record.course_id,
                    target_class_id=record.class_id,
                    draft_revision=review.draft_revision,
                    document_sha256=review.document_sha256,
                    submitted_by=review.submitted_by,
                    review_scope=review.scope,
                    review_status=review.status,
                )
            )
        return items, tuple(candidates)

    async def publish(self, command, materializer) -> PublishedVersionRecord:
        existing_id = self.ledger.publications_by_review.get(command.review_id)
        if existing_id is not None:
            return self.ledger.versions[existing_id].record
        review = self.ledger.reviews.get(command.review_id)
        current = await self.ledger.classroom(command.asset_id)
        if (
            review is None
            or current is None
            or review.tenant_id != command.tenant_id
            or current.tenant_id != command.tenant_id
            or review.asset_id != command.asset_id
            or review.asset_id != current.asset_id
            or review.draft_id != current.draft_id
            or review.scope != command.scope
            or review.class_id != command.class_id
            or (
                command.scope == "class"
                and command.class_id != current.class_id
            )
            or (command.scope != "class" and command.class_id is not None)
            or review.draft_revision != command.draft_revision
            or review.draft_revision != current.revision
            or current.classroom_version_id is None
            or current.validation_report is None
        ):
            raise ReviewConflict("publication review binding is stale")
        approved = review.status == "approved" and current.lifecycle_state == "approved"
        self_publish = (
            command.allow_self_publish
            and self.ledger.policy.teacher_self_publish
            and command.scope == "class"
            and command.class_id == current.class_id
            and current.owner_id == command.actor_id
            and review.submitted_by == command.actor_id
            and review.status == "pending"
            and current.lifecycle_state == "submitted"
        )
        if not approved and not self_publish:
            raise ReviewConflict("publication approval is unavailable")
        document = canonical_json_bytes(ClassroomDocument.model_validate(current.document))
        document_sha256 = hashlib.sha256(document).hexdigest()
        validation_report_sha256 = hashlib.sha256(
            canonical_json_bytes(current.validation_report)
        ).hexdigest()
        if (
            document_sha256 != review.document_sha256
            or document_sha256 != command.document_sha256
            or validation_report_sha256 != review.validation_report_sha256
        ):
            raise ReviewConflict("reviewed document binding is stale")
        version_number = 1 + sum(
            version.record.asset_id == current.asset_id
            for version in self.ledger.versions.values()
        )
        version_id = (
            "version-"
            + hashlib.sha256(
                f"{command.tenant_id}\0{command.review_id}".encode()
            ).hexdigest()[:24]
        )
        media_manifest_sha256 = publication_media_manifest_sha256(document)
        manifest_sha256 = hashlib.sha256(
            canonical_json_bytes(
                {
                    "tenantId": command.tenant_id,
                    "assetId": command.asset_id,
                    "reviewId": command.review_id,
                    "versionId": version_id,
                    "documentSha256": document_sha256,
                }
            )
        ).hexdigest()
        plan = PublicationMaterializationPlan(
            reservation_id=f"reservation-{command.review_id}",
            tenant_id=command.tenant_id,
            asset_id=command.asset_id,
            review_id=command.review_id,
            draft_id=current.draft_id,
            draft_revision=current.revision,
            source_version_id=current.classroom_version_id or "",
            version_id=version_id,
            version_number=version_number,
            document=document,
            document_sha256=document_sha256,
            validation_report_sha256=review.validation_report_sha256,
            media_manifest_sha256=media_manifest_sha256,
            manifest_sha256=manifest_sha256,
            media=(),
            status="prepared",
        )
        confirmed = await materializer.materialize(plan)
        if (
            confirmed.manifest_sha256 != plan.manifest_sha256
            or confirmed.media_manifest_sha256 != plan.media_manifest_sha256
            or confirmed.document.sha256 != plan.document_sha256
        ):
            raise ReviewConflict("publication materialization binding is invalid")
        published = PublishedVersionRecord(
            version_id=version_id,
            asset_id=current.asset_id,
            version_number=version_number,
            document_sha256=document_sha256,
            publication_scope=command.scope,
            class_id=command.class_id,
            idempotency_key=command.idempotency_key,
        )
        self.ledger.versions[version_id] = _LedgerVersion(
            record=published,
            tenant_id=current.tenant_id,
            owner_id=current.owner_id,
            course_id=current.course_id,
            class_id=current.class_id,
            title=current.title,
            review_id=review.id,
            draft_revision=current.revision,
            published_by=command.actor_id,
            document=bytes(document),
            media_manifest_sha256=media_manifest_sha256,
        )
        self.ledger.publications_by_review[review.id] = version_id
        lifecycle_state = "approved"
        if self_publish:
            lifecycle_state = transition("submitted", "approved")
        self.ledger.classrooms.records[current.asset_id] = replace(
            current,
            lifecycle_state=transition(lifecycle_state, "published"),
        )
        return published

    async def get_version_target(self, version_id: str) -> VersionTarget | None:
        version = self.ledger.versions.get(version_id)
        if version is None:
            return None
        return VersionTarget(
            tenant_id=version.tenant_id,
            version_id=version_id,
            asset_id=version.record.asset_id,
            course_id=version.course_id,
            publication_scope=version.record.publication_scope,
            publication_class_id=version.record.class_id,
        )

    async def assign(self, command) -> AssignmentRecord:
        existing = next(
            (
                assignment
                for assignment in self.ledger.assignments.values()
                if assignment.idempotency_key == command.idempotency_key
            ),
            None,
        )
        if existing is not None:
            return existing
        assignment = AssignmentRecord(
            assignment_id=f"assignment-{len(self.ledger.assignments) + 1}",
            tenant_id=command.tenant_id,
            asset_id=command.asset_id,
            version_id=command.version_id,
            class_id=command.class_id,
            assigned_by=command.actor_id,
            idempotency_key=command.idempotency_key,
            revoked_at=None,
        )
        self.ledger.assignments[assignment.assignment_id] = assignment
        return assignment

    async def get_assignment_target(self, assignment_id: str):
        del assignment_id
        return None

    async def get_migration(self, idempotency_key: str):
        del idempotency_key
        return None

    async def migrate(self, command):
        raise AssertionError(f"unexpected migration: {command}")


class _LedgerExportRepository:
    def __init__(self, ledger: _FlowLedger) -> None:
        self.ledger = ledger
        self.records: dict[str, ExportRecord] = {}
        self.version_reads: list[str] = []

    async def get_draft_source(self, asset_id: str):
        del asset_id
        return None

    async def get_version_source(self, version_id: str) -> ExportSource | None:
        self.version_reads.append(version_id)
        version = self.ledger.versions.get(version_id)
        if version is None:
            return None
        return ExportSource(
            tenant_id=version.tenant_id,
            asset_id=version.record.asset_id,
            owner_id=version.owner_id,
            course_id=version.course_id,
            class_id=version.class_id,
            classroom_draft_id=None,
            classroom_version_id=version_id,
            draft_revision=None,
            document=version.document,
            document_sha256=version.record.document_sha256,
            media_manifest_sha256=version.media_manifest_sha256,
            media=(),
        )

    async def reserve(self, command) -> ExportRecord:
        existing = self.records.get(command.idempotency_key)
        if existing is not None:
            return existing
        record = ExportRecord.from_command(command)
        self.records[command.idempotency_key] = record
        return record

    async def confirm_input(self, export_id, receipt) -> ExportRecord:
        for key, record in tuple(self.records.items()):
            if record.export_id == export_id:
                updated = replace(record, input_receipt=receipt)
                self.records[key] = updated
                return updated
        raise AssertionError("unknown export")

    async def bind_job(self, export_id: str, job_id: str) -> ExportRecord:
        for key, record in tuple(self.records.items()):
            if record.export_id == export_id:
                updated = replace(record, job_id=job_id, status="quota_reserved")
                self.records[key] = updated
                return updated
        raise AssertionError("unknown export")

    async def get(self, export_id: str) -> ExportRecord | None:
        return next(
            (record for record in self.records.values() if record.export_id == export_id),
            None,
        )


def _deterministic_zip(entries: dict[str, bytes]) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(
        buffer,
        "w",
        compression=zipfile.ZIP_DEFLATED,
        compresslevel=9,
    ) as archive:
        for name in sorted(entries):
            info = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o600 << 16
            archive.writestr(info, entries[name])
    return buffer.getvalue()


def _pptx_export(document_sha256: str) -> bytes:
    presentation = Presentation()
    presentation.core_properties.author = "yFeiSTAI E2E"
    presentation.core_properties.title = "Verified classroom export"
    presentation.core_properties.created = NOW.replace(tzinfo=None)
    presentation.core_properties.modified = NOW.replace(tzinfo=None)
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "Verified classroom export"
    slide.placeholders[1].text = f"Document SHA-256: {document_sha256}"
    raw = BytesIO()
    presentation.save(raw)
    with zipfile.ZipFile(BytesIO(raw.getvalue())) as archive:
        return _deterministic_zip(
            {name: archive.read(name) for name in archive.namelist()}
        )


class _CompletedExportJobs:
    _FORMATS = {
        "classroom_zip": ("classroom.zip", "application/zip"),
        "pptx": (
            "classroom.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        "offline_html": ("classroom.html", "text/html"),
    }

    def __init__(self, repository: _LedgerExportRepository, store: _BytesStore) -> None:
        self.repository = repository
        self.store = store
        self.commands = []

    def _artifact_body(self, command, record: ExportRecord) -> bytes:
        assert record.classroom_version_id is not None
        version = self.repository.ledger.versions[record.classroom_version_id]
        assert version.record.document_sha256 == command.document_sha256
        if command.export_format == "classroom_zip":
            manifest = canonical_json_bytes(
                {
                    "schemaVersion": 1,
                    "classroomVersionId": record.classroom_version_id,
                    "documentSha256": version.record.document_sha256,
                }
            )
            return _deterministic_zip(
                {
                    "classroom.json": version.document,
                    "manifest.json": manifest,
                }
            )
        if command.export_format == "pptx":
            return _pptx_export(version.record.document_sha256)
        assert command.export_format == "offline_html"
        return (
            '<?xml version="1.0" encoding="utf-8"?>'
            '<html xmlns="http://www.w3.org/1999/xhtml">'
            '<head><meta charset="utf-8"/><title>Verified classroom export</title></head>'
            f'<body data-document-sha256="{version.record.document_sha256}">'
            '<main><h1>Verified classroom export</h1></main></body></html>'
        ).encode()

    async def enqueue(self, command):
        self.commands.append(command)
        record = await self.repository.bind_job(command.export_id, command.job_id)
        relative_name, mime_type = self._FORMATS[command.export_format]
        body = self._artifact_body(command, record)
        object_key = (
            f"tenants/{command.tenant_id}/classroom-exports/"
            f"{command.export_id}/{relative_name}"
        )
        self.store.content[object_key] = body
        completed = replace(
            record,
            status="succeeded",
            progress_percent=100,
            relative_name=relative_name,
            object_key=object_key,
            sha256=hashlib.sha256(body).hexdigest(),
            size_bytes=len(body),
            mime_type=mime_type,
        )
        self.repository.records[command.idempotency_key] = completed
        return completed


async def _streaming_body(response) -> bytes:
    chunks = []
    async for chunk in response.body_iterator:
        chunks.append(chunk.encode() if isinstance(chunk, str) else chunk)
    return b"".join(chunks)


def _assert_valid_export_downloads(
    downloads: dict[str, bytes],
    document_sha256: str,
) -> None:
    with zipfile.ZipFile(BytesIO(downloads["classroom_zip"])) as archive:
        assert archive.testzip() is None
        assert set(archive.namelist()) == {"classroom.json", "manifest.json"}
        classroom_document = archive.read("classroom.json")
        assert hashlib.sha256(classroom_document).hexdigest() == document_sha256
        ClassroomDocument.model_validate_json(classroom_document)
        manifest = json.loads(archive.read("manifest.json"))
        assert manifest["documentSha256"] == document_sha256

    with zipfile.ZipFile(BytesIO(downloads["pptx"])) as archive:
        assert archive.testzip() is None
        assert {
            "[Content_Types].xml",
            "_rels/.rels",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/slides/slide1.xml",
        }.issubset(archive.namelist())
    presentation = Presentation(BytesIO(downloads["pptx"]))
    assert len(presentation.slides) == 1
    slide_text = "\n".join(
        shape.text
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text")
    )
    assert document_sha256 in slide_text

    html = ElementTree.fromstring(downloads["offline_html"])
    namespace = "{http://www.w3.org/1999/xhtml}"
    assert html.tag == f"{namespace}html"
    body = html.find(f"{namespace}body")
    assert body is not None
    assert body.attrib["data-document-sha256"] == document_sha256
    assert body.find(f"{namespace}main") is not None


async def _create_valid_grounded_classroom(
    classrooms: ClassroomService,
    repository: _FlowClassroomRepository,
    generation: _FlowGeneration,
    context: TenantContext,
    *,
    source_type: str,
    source_ref: str,
    idempotency_key: str,
    outline_title: str,
):
    request = classroom_request(
        content_mode="source_grounded",
        open_creation_acknowledged=False,
        source_type=source_type,
        source_ref=source_ref,
        requested_exports=["classroom_zip", "pptx", "offline_html"],
    )
    created = await classrooms.create(
        context,
        request,
        idempotency_key=idempotency_key,
    )
    assert created.outline is not None
    outline_payload = OutlineBundle.model_validate(created.outline).model_dump(
        mode="json",
        by_alias=True,
        exclude_none=True,
    )
    outline_payload["title"] = outline_title
    edited = await classrooms.update_outline(
        context,
        created.asset_id,
        outline_payload,
        expected_revision=created.revision,
    )
    assert edited.outline is not None
    confirmed = await classrooms.confirm_outline(
        context,
        created.asset_id,
        expected_revision=edited.revision,
        expected_outline_sha256=canonical_outline_sha256(
            OutlineBundle.model_validate(edited.outline)
        ),
    )
    await generation.complete(repository, created.asset_id)
    editing = await classrooms.get(context, created.asset_id)
    assert editing is not None
    assert editing.lifecycle_state == "editing"
    assert editing.status == "succeeded"
    validated = await classrooms.validate(context, created.asset_id)
    assert validated.validation_report is not None
    assert validated.validation_report["valid"] is True
    return validated, (
        created.lifecycle_state,
        confirmed.lifecycle_state,
        editing.lifecycle_state,
    )


async def _run_teacher_flow() -> dict[str, object]:
    teacher = _context("teacher-a", {"teacher"})
    reviewer = _context("reviewer-a", {"content_reviewer"})
    source_store = _BytesStore()
    pdf_body = b"%PDF-1.7 deterministic source"
    pdf_key = "tenants/tenant-a/sources/pdf-source-1/source.pdf"
    source_store.content[pdf_key] = pdf_body
    source_repository = _SourceRepository(
        pdf_key,
        hashlib.sha256(pdf_body).hexdigest(),
    )
    snapshots = SourceSnapshotBuilder(
        teacher,
        source_repository,
        knowledge_resolver=_knowledge_source,
        rag_service_factory=lambda _source: _RagService(),
        store_provider=_StoreProvider(source_store),
        pdf_extractor=_extract_pdf,
        clock=lambda: NOW,
    )
    classrooms_repository = _FlowClassroomRepository()
    generation = _FlowGeneration()
    classrooms = ClassroomService(
        classrooms_repository,
        TeachingBriefBuilder(teacher, snapshots),
        generation,
        None,
        clock=lambda: NOW,
    )
    validated, states = await _create_valid_grounded_classroom(
        classrooms,
        classrooms_repository,
        generation,
        teacher,
        source_type="knowledge_base",
        source_ref=KB_SOURCE_ID,
        idempotency_key="teacher-flow-kb",
        outline_title="Teacher-edited motion outline",
    )
    pdf_validated, pdf_states = await _create_valid_grounded_classroom(
        classrooms,
        classrooms_repository,
        generation,
        teacher,
        source_type="pdf",
        source_ref="binding-pdf",
        idempotency_key="teacher-flow-pdf",
        outline_title="Teacher-edited PDF outline",
    )
    assert tuple(item[0].source_type for item in source_repository.persisted) == (
        "knowledge_base",
        "pdf",
    )
    assert tuple(classrooms_repository.reconciled_assets) == (
        validated.asset_id,
        pdf_validated.asset_id,
    )

    edited_document = dict(validated.document)
    openmaic = edited_document["openmaic"]
    assert isinstance(openmaic, dict)
    scenes = openmaic["scenes"]
    assert isinstance(scenes, list) and isinstance(scenes[0], dict)
    scenes[0]["title"] = "Teacher-edited generated lesson"
    edited_document["fileSha256"] = "f" * 64
    edited = await classrooms.update_draft(
        teacher,
        validated.asset_id,
        edited_document,
        expected_revision=validated.revision,
    )
    assert edited.validation_report is None
    validated = await classrooms.validate(teacher, validated.asset_id)
    assert validated.validation_report is not None

    ledger = _FlowLedger(classrooms_repository)
    review_repository = _LedgerReviewRepository(ledger)
    reviews = ReviewService(review_repository)
    submitted = await reviews.submit(
        teacher,
        validated.asset_id,
        scope="class",
        class_id="class-a",
        idempotency_key="teacher-flow-review",
    )
    submitted_classroom = await classrooms.get(teacher, validated.asset_id)
    assert submitted_classroom is not None
    states = (*states, submitted_classroom.lifecycle_state)
    self_review_denied = False
    with pytest.raises(ReviewAccessDenied, match="self-review"):
        await reviews.approve(teacher, submitted.id, "self approval is forbidden")
    self_review_denied = True
    approved = await reviews.approve(reviewer, submitted.id, "reviewed by another user")
    approved_classroom = await classrooms.get(teacher, validated.asset_id)
    assert approved_classroom is not None
    states = (*states, approved_classroom.lifecycle_state)
    assert approved.reviewer_id == "reviewer-a"

    publications_repository = _LedgerPublicationRepository(ledger)
    publication_materializer = _LedgerPublicationMaterializer()
    publications = PublicationService(
        publications_repository,
        publication_materializer,
    )
    published = await publications.publish(
        teacher,
        validated.asset_id,
        scope="class",
        class_id="class-a",
        idempotency_key="teacher-flow-publish",
    )
    published_classroom = await classrooms.get(teacher, validated.asset_id)
    assert published_classroom is not None
    states = (*states, published_classroom.lifecycle_state)
    assignment = await publications.assign(
        teacher,
        published.version_id,
        class_id="class-a",
        idempotency_key="teacher-flow-assignment",
    )

    pdf_submitted = await reviews.submit(
        teacher,
        pdf_validated.asset_id,
        scope="class",
        class_id="class-a",
        idempotency_key="teacher-flow-pdf-review",
    )
    pdf_submitted_classroom = await classrooms.get(teacher, pdf_validated.asset_id)
    assert pdf_submitted_classroom is not None
    assert pdf_submitted_classroom.lifecycle_state == "submitted"
    assert pdf_submitted.draft_revision == pdf_validated.revision
    disabled_self_publish_denied = False
    with pytest.raises(PublicationAccessDenied):
        await publications.publish(
            teacher,
            pdf_validated.asset_id,
            scope="class",
            class_id="class-a",
            idempotency_key="teacher-flow-pdf-disabled-publish",
        )
    disabled_self_publish_denied = True

    ledger.policy = ReviewPolicy(teacher_self_publish=True)
    wrong_class_denied = False
    with pytest.raises(PublicationAccessDenied):
        await publications.publish(
            teacher,
            pdf_validated.asset_id,
            scope="class",
            class_id="class-b",
            idempotency_key="teacher-flow-pdf-wrong-class",
        )
    wrong_class_denied = True
    other_teacher = _context("teacher-b", {"teacher"})
    non_owner_denied = False
    with pytest.raises(PublicationAccessDenied):
        await publications.publish(
            other_teacher,
            pdf_validated.asset_id,
            scope="class",
            class_id="class-a",
            idempotency_key="teacher-flow-pdf-non-owner",
        )
    non_owner_denied = True
    self_published = await publications.publish(
        teacher,
        pdf_validated.asset_id,
        scope="class",
        class_id="class-a",
        idempotency_key="teacher-flow-pdf-self-publish",
    )
    pdf_published_classroom = await classrooms.get(teacher, pdf_validated.asset_id)
    assert pdf_published_classroom is not None
    pdf_version = ledger.versions[self_published.version_id]
    pdf_plan = next(
        plan
        for plan in publication_materializer.plans
        if plan.review_id == pdf_submitted.id
    )
    self_published_immutable = (
        pdf_published_classroom.lifecycle_state == "published"
        and pdf_version.review_id == pdf_submitted.id
        and pdf_version.draft_revision == pdf_submitted.draft_revision
        and pdf_plan.draft_revision == pdf_submitted.draft_revision
        and pdf_plan.document_sha256 == pdf_submitted.document_sha256
        and hashlib.sha256(pdf_version.document).hexdigest()
        == self_published.document_sha256
        == pdf_submitted.document_sha256
    )
    with pytest.raises(InvalidClassroomTransition):
        transition("published", "editing")

    exports_repository = _LedgerExportRepository(ledger)
    export_source = await exports_repository.get_version_source(published.version_id)
    assert export_source is not None
    immutable_version = (
        hashlib.sha256(export_source.document).hexdigest()
        == published.document_sha256
        == ledger.versions[published.version_id].record.document_sha256
    )
    export_store = _BytesStore()
    export_materializer = ExportMemoryMaterializer()
    export_jobs = _CompletedExportJobs(exports_repository, export_store)
    exports = ClassroomExportService(
        exports_repository,
        export_materializer,
        export_jobs,
        mp4_enabled=lambda _tenant_id: False,
    )
    downloads: dict[str, bytes] = {}
    for export_format in ("classroom_zip", "pptx", "offline_html"):
        export = await exports.create_for_version(
            teacher,
            published.version_id,
            export_format,
            idempotency_key=f"teacher-flow-{export_format}",
        )
        assert export.status == "succeeded"
        assert export.input_document_sha256 == published.document_sha256
        response = await download_classroom_export(
            export.export_id,
            context=teacher,
            service=exports,
            stores=_StoreProvider(export_store),
        )
        assert "tenants/" not in response.headers["content-disposition"]
        downloads[export_format] = await _streaming_body(response)
    assert len(export_materializer.plans) == 3
    assert all(
        plan.classroom_version_id == published.version_id
        and plan.document_sha256 == published.document_sha256
        for plan in export_materializer.plans
    )

    return {
        "sourceKinds": tuple(item[0].source_type for item in source_repository.persisted),
        "states": tuple(states),
        "selfReviewDenied": self_review_denied,
        "publishedVersionId": published.version_id,
        "assignmentVersionId": assignment.version_id,
        "immutableVersion": immutable_version,
        "selfPublishPolicy": {
            "disabledDenied": disabled_self_publish_denied,
            "wrongClassDenied": wrong_class_denied,
            "nonOwnerDenied": non_owner_denied,
            "ownerPublishedVersionId": self_published.version_id,
            "immutableRevisionAndSha": self_published_immutable,
        },
        "reviewDerivedPublication": (
            publication_materializer.plans[0].review_id == approved.id
            and ledger.versions[published.version_id].review_id == approved.id
        ),
        "publishedDocumentSha256": published.document_sha256,
        "exportDocumentSha256": export_source.document_sha256,
        "exportSourceVersionId": export_source.classroom_version_id,
        "pdfReachedValidEditing": (
            pdf_states[-1] == "editing"
            and pdf_validated.lifecycle_state == "editing"
            and pdf_validated.validation_report is not None
            and pdf_validated.validation_report["valid"] is True
        ),
        "downloads": downloads,
    }


@pytest.mark.asyncio
async def test_teacher_classroom_flow_reaches_controlled_exports() -> None:
    evidence = await _run_teacher_flow()

    assert evidence["sourceKinds"] == ("knowledge_base", "pdf")
    assert evidence["states"] == (
        "awaiting_outline",
        "generating_content",
        "editing",
        "submitted",
        "approved",
        "published",
    )
    assert evidence["selfReviewDenied"] is True
    assert evidence["assignmentVersionId"] == evidence["publishedVersionId"]
    assert evidence["immutableVersion"] is True
    assert evidence["selfPublishPolicy"]["disabledDenied"] is True
    assert evidence["selfPublishPolicy"]["wrongClassDenied"] is True
    assert evidence["selfPublishPolicy"]["nonOwnerDenied"] is True
    assert evidence["selfPublishPolicy"]["immutableRevisionAndSha"] is True
    assert (
        evidence["selfPublishPolicy"]["ownerPublishedVersionId"]
        != evidence["publishedVersionId"]
    )
    assert evidence["reviewDerivedPublication"] is True
    assert evidence["publishedDocumentSha256"] == evidence["exportDocumentSha256"]
    assert evidence["exportSourceVersionId"] == evidence["publishedVersionId"]
    assert evidence["pdfReachedValidEditing"] is True
    assert set(evidence["downloads"]) == {"classroom_zip", "pptx", "offline_html"}
    _assert_valid_export_downloads(
        evidence["downloads"],
        evidence["publishedDocumentSha256"],
    )
